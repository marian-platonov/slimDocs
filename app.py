"""SlimDocs - Extract and reformat file content for AI workflows."""

import io
import json
import os
import re
import tarfile
import tempfile
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from datetime import datetime

import pandas as pd
import streamlit as st

# ── wevtapi.dll - Windows-native EVTX parser (superior to evtx library) ──────
# Loaded once at startup; None on non-Windows or if the DLL is unavailable.
import ctypes
import sys as _sys

_wevt: "ctypes.WinDLL | None" = None
if _sys.platform == "win32":
    try:
        _wevt = ctypes.WinDLL("wevtapi.dll", use_last_error=True)
        _wevt.EvtQuery.restype  = ctypes.c_void_p
        _wevt.EvtQuery.argtypes = [
            ctypes.c_void_p, ctypes.c_wchar_p, ctypes.c_wchar_p, ctypes.c_uint,
        ]
        _wevt.EvtNext.restype   = ctypes.c_bool
        _wevt.EvtNext.argtypes  = [
            ctypes.c_void_p, ctypes.c_uint,
            ctypes.POINTER(ctypes.c_void_p),
            ctypes.c_uint, ctypes.c_uint, ctypes.POINTER(ctypes.c_uint),
        ]
        _wevt.EvtRender.restype  = ctypes.c_bool
        _wevt.EvtRender.argtypes = [
            ctypes.c_void_p, ctypes.c_void_p, ctypes.c_uint, ctypes.c_uint,
            ctypes.c_void_p, ctypes.POINTER(ctypes.c_uint), ctypes.POINTER(ctypes.c_uint),
        ]
        _wevt.EvtClose.restype   = ctypes.c_bool
        _wevt.EvtClose.argtypes  = [ctypes.c_void_p]
        _wevt.EvtFormatMessage.restype  = ctypes.c_bool
        _wevt.EvtFormatMessage.argtypes = [
            ctypes.c_void_p, ctypes.c_void_p, ctypes.c_uint, ctypes.c_uint,
            ctypes.c_void_p, ctypes.c_uint, ctypes.c_uint,
            ctypes.c_void_p, ctypes.POINTER(ctypes.c_uint),
        ]
    except Exception:
        _wevt = None

_EVT_QUERY_FILE_PATH      = 0x2
_EVT_RENDER_EVENT_XML     = 0x1
_EVT_FORMAT_MESSAGE_EVENT = 0x1
_INFINITE                 = 0xFFFFFFFF
_WEVT_BATCH               = 100


def _wevt_format_message(h: int) -> str:
    try:
        used = ctypes.c_uint(0)
        _wevt.EvtFormatMessage(
            None, h, 0, 0, None, _EVT_FORMAT_MESSAGE_EVENT, 0, None, ctypes.byref(used),
        )
        if used.value == 0:
            return ""
        buf = ctypes.create_unicode_buffer(used.value)
        ok = _wevt.EvtFormatMessage(
            None, h, 0, 0, None, _EVT_FORMAT_MESSAGE_EVENT,
            used.value, buf, ctypes.byref(used),
        )
        return buf.value.strip() if ok else ""
    except Exception:
        return ""


def _wevt_render_xml(h: int) -> str:
    buf_used   = ctypes.c_uint(0)
    prop_count = ctypes.c_uint(0)
    _wevt.EvtRender(None, h, _EVT_RENDER_EVENT_XML, 0, None,
                    ctypes.byref(buf_used), ctypes.byref(prop_count))
    size = buf_used.value
    if size == 0:
        raise RuntimeError("EvtRender size probe returned 0")
    buf = ctypes.create_unicode_buffer(size // 2 + 4)
    ok  = _wevt.EvtRender(None, h, _EVT_RENDER_EVENT_XML, size, buf,
                          ctypes.byref(buf_used), ctypes.byref(prop_count))
    if not ok:
        raise OSError(ctypes.get_last_error(), "EvtRender")
    return buf.value


def _wevt_records(filepath: str):
    """Generator: yields (xml_str, message_text) per EVTX record via wevtapi."""
    query = _wevt.EvtQuery(None, filepath, "*", _EVT_QUERY_FILE_PATH)
    if not query:
        raise OSError(ctypes.get_last_error(), "EvtQuery")
    arr      = (ctypes.c_void_p * _WEVT_BATCH)()
    returned = ctypes.c_uint(0)
    try:
        while True:
            returned.value = 0
            ok = _wevt.EvtNext(query, _WEVT_BATCH, arr, _INFINITE, 0, ctypes.byref(returned))
            n  = returned.value
            if n == 0:
                break
            for i in range(n):
                h = arr[i]
                try:
                    yield _wevt_render_xml(h), _wevt_format_message(h)
                finally:
                    _wevt.EvtClose(h)
            if not ok:
                break
    finally:
        _wevt.EvtClose(query)


# ── Page config ──────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="SlimDocs",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Extension registry (mirrors fileScopeGUI scanner categories) ─────────────
_EXT_CATEGORIES = {
    "Documents":    {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".odt", ".rtf", ".md", ".txt"},
    "Spreadsheets": {".xlsx", ".xls", ".csv", ".tsv", ".ods"},
    "Web":          {".html", ".htm"},
    "Archives":     {".zip", ".tar", ".gz", ".bz2", ".rar", ".xz", ".7z"},
    "Images":       {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif", ".webp"},
    "Logs":         {".log", ".out", ".err", ".evtx"},
    "Code":         {
        ".py", ".pyw", ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx",
        ".java", ".kt", ".scala", ".groovy", ".cpp", ".cxx", ".cc",
        ".c", ".h", ".hpp", ".hxx", ".cs", ".vb", ".go", ".rs",
        ".rb", ".rake", ".php", ".swift", ".r", ".rmd", ".m", ".mm",
        ".lua", ".pl", ".pm", ".sql", ".ddl", ".dml",
        ".sh", ".bash", ".zsh", ".fish", ".bat", ".cmd", ".ps1", ".psm1",
        ".asm", ".s", ".dart", ".elm", ".ex", ".exs",
        ".clj", ".cljs", ".erl", ".hs",
    },
    "Config/Data":  {
        ".json", ".jsonc", ".json5", ".yaml", ".yml", ".toml",
        ".xml", ".ini", ".cfg", ".conf", ".config", ".env",
        ".properties", ".plist", ".editorconfig", ".gitignore",
        ".gitattributes", ".dockerfile", ".tf", ".tfvars",
        ".gradle", ".maven",
    },
    "Videos":       {".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm",
                     ".m4v", ".mpg", ".mpeg", ".3gp"},
}

_SUPPORTED_EXTENSIONS: set = {ext for exts in _EXT_CATEGORIES.values() for ext in exts}

# Extensions inside archives that get their text extracted
_ARCHIVE_TEXT_EXTS = {
    ".txt", ".log", ".md", ".csv", ".tsv", ".html", ".xml",
    ".json", ".jsonc", ".yaml", ".yml", ".ini", ".conf", ".cfg",
    ".sh", ".bat", ".ps1", ".py", ".js", ".ts", ".sql", ".evtx",
}

# Extensions treated as plain text (read as-is)
_PLAINTEXT_EXTS = (
    _EXT_CATEGORIES["Code"]
    | _EXT_CATEGORIES["Config/Data"]
    | _EXT_CATEGORIES["Logs"]
    | {".md", ".txt", ".rtf"}
)

_FORMAT_OPTIONS = {
    "md - Markdown": "md",
    "txt - Plain text": "txt",
    "json - Chunked Claude-ready JSON": "json",
    "duckdb - DuckDB database": "duckdb",
}

_DOWNLOADS_DIR = Path.home() / "Downloads"
_PREVIEW_LIMIT = 20_000

# Populated by _extract_pdf_images_ocr; consumed + cleared in process_files.
# Maps str(source_pdf_path) → list of (relative_output_name, raw_bytes)
_pdf_img_cache: dict = {}


# ── Session state ────────────────────────────────────────────────────────────
def _init_session():
    defaults = {
        "input_mode": "Single File",
        "input_path": "",
        "urls_text":  "",
        "output_dir": str(_DOWNLOADS_DIR),
        "results":    None,
        "processing": False,
        "content_map": {},
        "logs": [],
        "session_history": {},
        "selected_session_ts": None,
        "show_nav_hint": False,
        "session_totals": {
            "runs": 0, "files_ok": 0, "errors": 0,
            "tokens_before": 0, "tokens_after": 0, "tokens_saved": 0,
        },
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v


# ── Low-level byte decoder ────────────────────────────────────────────────────
def _decode_bytes(data: bytes) -> str:
    for enc in ("utf-8-sig", "utf-16", "utf-8", "latin-1"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, ValueError):
            continue
    return data.decode("latin-1", errors="replace")


# ── Extractor functions ───────────────────────────────────────────────────────

def _extract_plaintext(path: Path) -> str:
    return _decode_bytes(path.read_bytes())


def _extract_csv(path: Path) -> str:
    import csv
    rows = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f):
            rows.append("\t".join(row))
    return "\n".join(rows)


def _extract_tsv(path: Path) -> str:
    import csv
    rows = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f, delimiter="\t"):
            rows.append("\t".join(row))
    return "\n".join(rows)


def _extract_pdf(path: Path) -> str:
    """Extract text + tables from every page; supplement with embedded-image OCR."""
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber not installed - run: pip install pdfplumber")

    pages_out = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            parts = [f"## Page {i}"]

            # ── Tables (structured Markdown) ──────────────────────────────
            try:
                tables = page.extract_tables()
                for t_idx, table in enumerate(tables, 1):
                    if not table:
                        continue
                    md_rows = []
                    for j, row in enumerate(table):
                        cells = [str(c or "").strip().replace("\n", " ") for c in row]
                        md_rows.append("| " + " | ".join(cells) + " |")
                        if j == 0:
                            md_rows.append("|" + "|".join(" --- " for _ in row) + "|")
                    parts.append(f"### Table {t_idx}\n" + "\n".join(md_rows))
            except Exception:
                pass

            # ── Visible text ──────────────────────────────────────────────
            text = page.extract_text()
            if text and text.strip():
                parts.append(text.strip())

            # ── Note embedded images/charts ───────────────────────────────
            if page.images:
                parts.append(
                    f"*[{len(page.images)} image/chart object(s) on this page - "
                    f"see Embedded Image OCR section below]*"
                )

            if len(parts) > 1:
                pages_out.append("\n\n".join(parts))

    result = "\n\n".join(pages_out)

    # ── Supplement: OCR every embedded image via PyMuPDF ─────────────────
    ocr_supplement = _extract_pdf_images_ocr(path)
    if ocr_supplement:
        result += "\n\n---\n\n## Embedded Image OCR\n\n" + ocr_supplement

    return result


def _extract_pdf_images_ocr(path: Path) -> str:
    """
    Pull every embedded raster image from a PDF via PyMuPDF, save the raw
    bytes to _pdf_img_cache (consumed later by process_files), embed a
    Markdown image reference, and append any OCR text if Tesseract is
    available.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ""
    try:
        from PIL import Image as _PILImage
    except ImportError:
        return ""

    try:
        import pytesseract
        _pytesseract_ok = True
    except ImportError:
        _pytesseract_ok = False

    def _ocr(img):
        if not _pytesseract_ok:
            return None
        try:
            return pytesseract.image_to_string(img).strip()
        except pytesseract.TesseractNotFoundError:
            return "__TESSERACT_MISSING__"
        except Exception:
            return ""

    stem = path.stem
    cache_key = str(path)
    _pdf_img_cache[cache_key] = []

    doc = fitz.open(str(path))
    seen_xrefs: set = set()
    ocr_parts: list = []
    tesseract_missing_reported = False
    img_counter = 0

    for page_num, page in enumerate(doc, 1):
        page_parts: list = []

        for img_ref in page.get_images(full=True):
            xref = img_ref[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                base_img = doc.extract_image(xref)
                raw_bytes = base_img["image"]
                img_ext   = base_img.get("ext", "png")
                img = _PILImage.open(io.BytesIO(raw_bytes)).convert("RGB")

                # Skip tiny decorative images (icons, bullets, separators …)
                if img.width < 50 or img.height < 50:
                    continue

                img_counter += 1
                # relative path used both in the cache and as the Markdown ref
                rel_name = f"{stem}_images/page{page_num}_img{img_counter}.{img_ext}"
                _pdf_img_cache[cache_key].append((rel_name, raw_bytes))

                header = (
                    f"### Embedded image - page {page_num}, "
                    f"#{img_counter} ({img.width}×{img.height} px)\n"
                    f"![image]({rel_name})"
                )

                ocr_text = _ocr(img)
                if ocr_text == "__TESSERACT_MISSING__":
                    if not tesseract_missing_reported:
                        tesseract_missing_reported = True
                        page_parts.append(
                            "**OCR unavailable** - Tesseract binary not found.  \n"
                            "Install from https://github.com/UB-Mannheim/tesseract/wiki "
                            "and add its directory to your PATH, then re-run extraction."
                        )
                    page_parts.append(f"{header}  \n*[OCR skipped - Tesseract not installed]*")
                elif ocr_text:
                    page_parts.append(f"{header}  \n**OCR text:**\n{ocr_text}")
                else:
                    page_parts.append(f"{header}  \n*[No text detected by OCR]*")
            except Exception:
                pass

        ocr_parts.extend(page_parts)

    return "\n\n".join(ocr_parts)


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed - run: pip install python-docx")
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            prefix = "## " if para.style.name.startswith("Heading") else ""
            lines.append(f"{prefix}{para.text}")
    return "\n".join(lines)


def _extract_odt(path: Path) -> str:
    """OpenDocument Text - extract content.xml from the ZIP container."""
    try:
        from defusedxml import ElementTree as DET
    except ImportError:
        raise RuntimeError("defusedxml not installed - run: pip install defusedxml")
    with zipfile.ZipFile(path, "r") as zf:
        if "content.xml" not in zf.namelist():
            raise ValueError("content.xml not found inside ODT file")
        raw = zf.read("content.xml")
    root = DET.fromstring(raw)
    parts = []
    for el in root.iter():
        tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        if tag in ("p", "h"):
            text = "".join(el.itertext()).strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _extract_rtf(path: Path) -> str:
    """Strip RTF control sequences and return readable text."""
    raw = _decode_bytes(path.read_bytes())
    text = re.sub(r"\\[a-zA-Z]+\-?[0-9]* ?", " ", raw)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"[{}\\]", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl not installed - run: pip install openpyxl")
    wb = openpyxl.load_workbook(path, data_only=True)
    sections = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                rows.append("\t".join("" if c is None else str(c) for c in row))
        if rows:
            sections.append(f"## Sheet: {name}\n" + "\n".join(rows))
    return "\n\n".join(sections)


def _extract_ods(path: Path) -> str:
    try:
        df = pd.read_excel(path, engine="odf")
        return df.to_string(index=False)
    except Exception as e:
        if "odf" in str(e).lower() or "odfpy" in str(e).lower():
            raise RuntimeError("odfpy not installed - run: pip install odfpy")
        raise


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed - run: pip install python-pptx")
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        if len(lines) > 1:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _extract_image(path: Path) -> str:
    try:
        from PIL import Image, ExifTags
        import pytesseract
    except ImportError:
        raise RuntimeError(
            "pillow/pytesseract not installed - run: pip install pillow pytesseract\n"
            "Also install Tesseract: https://github.com/UB-Mannheim/tesseract/wiki"
        )
    with Image.open(path) as img:
        lines = [
            f"Image: {path.name}",
            f"Format: {img.format or path.suffix.upper().lstrip('.')}",
            f"Dimensions: {img.size[0]}x{img.size[1]} px",
            f"Mode: {img.mode}",
        ]
        try:
            exif = img._getexif()  # type: ignore[attr-defined]
            if exif:
                for tag_id, value in exif.items():
                    tag = ExifTags.TAGS.get(tag_id, "")
                    if tag in ("Make", "Model", "DateTime"):
                        lines.append(f"{tag}: {value}")
        except Exception:
            pass
        try:
            ocr = pytesseract.image_to_string(img).strip()
            if ocr:
                lines.append(f"\nOCR Text:\n{ocr}")
        except pytesseract.TesseractNotFoundError:
            lines.append(
                "\nOCR skipped: Tesseract binary not found.\n"
                "Install it from https://github.com/UB-Mannheim/tesseract/wiki "
                "and make sure it is added to your PATH."
            )
        except Exception:
            pass
    return "\n".join(lines)


def _extract_html(path: Path) -> str:
    try:
        import trafilatura
    except ImportError:
        raise RuntimeError("trafilatura not installed - run: pip install trafilatura")
    html = _decode_bytes(path.read_bytes())
    result = trafilatura.extract(
        html, output_format="markdown",
        include_links=True, include_images=False, favor_recall=True,
    )
    return result or ""


def _extract_archive(path: Path) -> str:
    """List archive members and extract text from supported inner files."""
    ext = path.suffix.lower()
    lines = [f"Archive: {path.name}", "=" * 60]
    text_sections: list = []
    _extracted_count = 0
    _skipped_count   = 0

    def _add_text(member_name: str, data: bytes):
        nonlocal _extracted_count
        text = _decode_bytes(data).strip()
        if text:
            text_sections.append(f"=== {member_name} ===\n{text}")
            _extracted_count += 1

    _ARCHIVE_EXTS = {".zip", ".tar", ".gz", ".bz2", ".rar", ".xz", ".7z"}

    def _try_extract_file(name: str, data: bytes):
        """Try text extraction, then rich extractor; update counters."""
        nonlocal _skipped_count, _extracted_count
        inner_ext = Path(name).suffix.lower()
        if inner_ext in _ARCHIVE_TEXT_EXTS:
            _add_text(name, data)
        elif inner_ext not in _ARCHIVE_EXTS:
            extractor = _FILE_EXTRACTORS.get(inner_ext)
            if extractor:
                try:
                    tmp_dir  = Path(tempfile.mkdtemp())
                    tmp_path = tmp_dir / Path(name).name
                    tmp_path.write_bytes(data)
                    text = extractor(tmp_path)
                    text_sections.append(f"=== {name} ===\n{text}")
                    _extracted_count += 1
                    imgs = _pdf_img_cache.pop(str(tmp_path), [])
                    if imgs:
                        _pdf_img_cache.setdefault(str(path), []).extend(imgs)
                except Exception as exc:
                    lines.append(f"  [extract error: {exc}]")
                finally:
                    try:
                        tmp_path.unlink(missing_ok=True)
                        tmp_dir.rmdir()
                    except Exception:
                        pass
            else:
                _skipped_count += 1
        else:
            _skipped_count += 1  # nested archive - skip to avoid recursion

    if ext == ".zip":
        with zipfile.ZipFile(path, "r") as zf:
            for entry in zf.infolist():
                tag = "[DIR] " if entry.is_dir() else "[FILE]"
                lines.append(f"{tag} {entry.filename}  ({entry.file_size:,} B)")
                if not entry.is_dir() and entry.file_size > 0:
                    try:
                        _try_extract_file(entry.filename, zf.read(entry))
                    except Exception as e:
                        lines.append(f"  [read error: {e}]")

    elif ext in {".tar", ".gz", ".bz2", ".xz"}:
        with tarfile.open(path, "r:*") as tf:
            for member in tf.getmembers():
                tag = "[DIR] " if member.isdir() else "[FILE]"
                lines.append(f"{tag} {member.name}  ({member.size:,} B)")
                if member.isfile() and member.size > 0:
                    try:
                        f = tf.extractfile(member)
                        if f:
                            _try_extract_file(member.name, f.read())
                    except Exception as e:
                        lines.append(f"  [read error: {e}]")

    elif ext == ".rar":
        try:
            import rarfile
        except ImportError:
            raise RuntimeError("rarfile not installed - run: pip install rarfile")
        with rarfile.RarFile(path, "r") as rf:
            for info in rf.infolist():
                tag = "[DIR] " if info.isdir() else "[FILE]"
                lines.append(f"{tag} {info.filename}  ({info.file_size:,} B)")
                if not info.isdir():
                    try:
                        _try_extract_file(info.filename, rf.read(info.filename))
                    except Exception as e:
                        lines.append(f"  [read error: {e}]")

    elif ext == ".7z":
        try:
            import py7zr
        except ImportError:
            raise RuntimeError("py7zr not installed - run: pip install py7zr")
        with py7zr.SevenZipFile(path, mode="r") as sz:
            all_files = sz.getnames()
            for name in all_files:
                lines.append(f"[FILE] {name}")
            extracted = sz.read() or {}
            for name, bio in extracted.items():
                if bio is None:
                    continue
                try:
                    _try_extract_file(name, bio.read())
                except Exception as e:
                    lines.append(f"  [read error: {e}]")
    else:
        raise ValueError(f"Unsupported archive format: {ext}")

    summary = (
        f"\n--- Extraction summary: {_extracted_count} file(s) extracted"
        + (f", {_skipped_count} skipped (binary / unsupported extension)" if _skipped_count else "")
        + " ---"
    )
    lines.append(summary)

    result = "\n".join(lines)
    if text_sections:
        result += "\n\n" + "\n\n".join(text_sections)
    return result


def _evtx_parse_xml(xml_str: str, msg: str = "") -> str | None:
    """
    Parse one EVTX XML record and return a formatted log line.
    Shared by both the wevtapi and evtx-library paths.
    `msg` is the human-readable EvtFormatMessage text (wevtapi only).
    Returns None when the record should be skipped.
    """
    try:
        from defusedxml import ElementTree as DET
    except ImportError:
        raise RuntimeError("defusedxml not installed - run: pip install defusedxml")
    _NS = "http://schemas.microsoft.com/win/2004/08/events/event"
    _LEVELS = {
        "0": "Info", "1": "Critical", "2": "Error",
        "3": "Warning", "4": "Information", "5": "Verbose",
    }
    root   = DET.fromstring(xml_str)
    ns     = {"e": _NS}
    system = root.find("e:System", ns)
    if system is None:
        return None

    def _txt(tag: str) -> str:
        el = system.find(f"e:{tag}", ns)
        return (el.text or "") if el is not None else ""

    level_str   = _LEVELS.get(_txt("Level"), f"Level{_txt('Level')}")
    provider_el = system.find("e:Provider", ns)
    source      = provider_el.get("Name", "") if provider_el is not None else ""
    time_el     = system.find("e:TimeCreated", ns)
    timestamp   = time_el.get("SystemTime", "") if time_el is not None else ""

    # EventData
    data_parts: list = []
    event_data = root.find("e:EventData", ns)
    if event_data is not None:
        for el in event_data:
            val = " ".join((el.text or "").splitlines()).strip()
            if val:
                name_attr = el.get("Name", "")
                data_parts.append(f"{name_attr}: {val}" if name_attr else val)

    # UserData (fallback when EventData is absent)
    user_data = root.find("e:UserData", ns)
    if user_data is not None and not data_parts:
        for child in user_data:
            for el in child:
                val = " ".join((el.text or "").splitlines()).strip()
                if val:
                    data_parts.append(f"{el.tag.split('}')[-1]}: {val}")

    line = (
        f"[{timestamp}] [{level_str}] "
        f"EventID={_txt('EventID')} Source={source} "
        f"Computer={_txt('Computer')} Channel={_txt('Channel')}"
    )
    if data_parts:
        line += "\n  " + " | ".join(data_parts[:15])
    if msg:
        line += "\n  Message: " + " ".join(msg.splitlines())[:400]

    return line


def _extract_evtx(path: Path) -> str:
    """
    Windows Event Log (.evtx) extraction.

    Strategy (in order):
      1. wevtapi.dll (Windows-native) - includes EvtFormatMessage human-readable
         descriptions, the richest output possible.
      2. evtx Python library - cross-platform fallback.
    """
    _LEVELS = {
        "0": "Info", "1": "Critical", "2": "Error",
        "3": "Warning", "4": "Information", "5": "Verbose",
    }
    lines: list        = []
    level_counts: dict = {}
    parse_errors       = 0

    # ── Path 1: wevtapi.dll ───────────────────────────────────────────────
    if _wevt is not None:
        try:
            for xml_str, msg in _wevt_records(str(path)):
                try:
                    line = _evtx_parse_xml(xml_str, msg)
                    if line is None:
                        continue
                    # tally level from the line header for the summary
                    level_tag = line.split("] [")[1].split("]")[0] if "] [" in line else "?"
                    level_counts[level_tag] = level_counts.get(level_tag, 0) + 1
                    lines.append(line)
                except (ET.ParseError, Exception):
                    parse_errors += 1
        except OSError as e:
            # wevtapi failed (e.g. file locked, corrupt) - fall through to evtx
            lines = []
            level_counts = {}
            parse_errors = 0

    # ── Path 2: evtx library fallback ────────────────────────────────────
    if not lines:
        try:
            from evtx import PyEvtxParser
        except ImportError:
            raise RuntimeError(
                "evtx library not installed - run: pip install evtx\n"
                "(wevtapi.dll is also unavailable on this system)"
            )
        try:
            parser = PyEvtxParser(io.BytesIO(path.read_bytes()))
        except Exception as e:
            raise RuntimeError(f"Could not open EVTX file: {e}")

        for record in parser.records():
            if isinstance(record, RuntimeError):
                parse_errors += 1
                continue
            try:
                xml_str = record.get("data", "")
                if not xml_str:
                    continue
                # Patch missing timestamp from record metadata
                ts_fallback = str(record.get("timestamp", ""))
                if ts_fallback and "TimeCreated" not in xml_str:
                    pass  # evtx library embeds the timestamp in the XML already
                line = _evtx_parse_xml(xml_str)
                if line is None:
                    continue
                level_tag = line.split("] [")[1].split("]")[0] if "] [" in line else "?"
                level_counts[level_tag] = level_counts.get(level_tag, 0) + 1
                lines.append(line)
            except (ET.ParseError, Exception):
                parse_errors += 1

    # ── Summary header ────────────────────────────────────────────────────
    header_parts = [f"EVTX: {path.name}", f"Total events: {len(lines)}"]
    if level_counts:
        header_parts.append(
            "Levels: " + "  ".join(f"{k}={v}" for k, v in sorted(level_counts.items()))
        )
    if parse_errors:
        header_parts.append(f"Parse errors: {parse_errors}")
    method = "wevtapi.dll" if (_wevt is not None and lines) else "evtx library"
    header_parts.append(f"Parser: {method}")

    header = "\n".join(header_parts) + "\n" + "=" * 70
    return header + "\n\n" + "\n\n".join(lines)


def _extract_video(path: Path) -> str:
    size_mb = path.stat().st_size / 1_048_576
    return (
        f"Video file: {path.name}\n"
        f"Format: {path.suffix.upper().lstrip('.')}\n"
        f"File size: {size_mb:.2f} MB\n"
        "(Content extraction not supported for video - file metadata only.)"
    )


# ── File extractor dispatch table ─────────────────────────────────────────────
_FILE_EXTRACTORS: dict = {}

# Documents
for _e in (".pdf",):                    _FILE_EXTRACTORS[_e] = _extract_pdf
for _e in (".docx", ".doc"):            _FILE_EXTRACTORS[_e] = _extract_docx
for _e in (".odt",):                    _FILE_EXTRACTORS[_e] = _extract_odt
for _e in (".rtf",):                    _FILE_EXTRACTORS[_e] = _extract_rtf
for _e in (".pptx", ".ppt"):            _FILE_EXTRACTORS[_e] = _extract_pptx
# Spreadsheets
for _e in (".xlsx", ".xls"):            _FILE_EXTRACTORS[_e] = _extract_xlsx
for _e in (".ods",):                    _FILE_EXTRACTORS[_e] = _extract_ods
for _e in (".csv",):                    _FILE_EXTRACTORS[_e] = _extract_csv
for _e in (".tsv",):                    _FILE_EXTRACTORS[_e] = _extract_tsv
# Web
for _e in (".html", ".htm"):            _FILE_EXTRACTORS[_e] = _extract_html
# Archives
for _e in (".zip", ".tar", ".gz",
           ".bz2", ".rar", ".xz",
           ".7z"):                       _FILE_EXTRACTORS[_e] = _extract_archive
# Images
for _e in (".png", ".jpg", ".jpeg",
           ".gif", ".bmp", ".tiff",
           ".tif", ".webp"):            _FILE_EXTRACTORS[_e] = _extract_image
# Logs
for _e in (".evtx",):                   _FILE_EXTRACTORS[_e] = _extract_evtx
# Videos
for _e in _EXT_CATEGORIES["Videos"]:   _FILE_EXTRACTORS[_e] = _extract_video
# Plain text: code, config, logs (.log/.out/.err), .md, .txt, .rtf already done
for _e in _PLAINTEXT_EXTS:
    if _e not in _FILE_EXTRACTORS:
        _FILE_EXTRACTORS[_e] = _extract_plaintext


def extract_from_file(path: Path) -> str:
    extractor = _FILE_EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        raise ValueError(f"Unsupported file type: '{path.suffix}'")
    return extractor(path)


# ── URL fetching ──────────────────────────────────────────────────────────────

def _fetch_confluence(url: str) -> str:
    import requests
    from urllib.parse import urlparse
    parsed = urlparse(url)
    match = re.search(r"/pages/(\d+)(?:/|$)", parsed.path)
    if not match:
        raise ValueError(f"Cannot extract Confluence page ID from: {url}")
    page_id = match.group(1)
    api_url = (
        f"{parsed.scheme}://{parsed.netloc}"
        f"/wiki/rest/api/content/{page_id}?expand=body.view"
    )
    email = os.environ.get("ATLASSIAN_EMAIL", "")
    token = os.environ.get("ATLASSIAN_API_TOKEN", "")
    auth = (email, token) if email and token else None
    resp = requests.get(
        api_url, timeout=15, auth=auth,
        headers={"User-Agent": "Mozilla/5.0 (compatible; SlimDocs/1.0)"},
    )
    if resp.status_code in (401, 403):
        if not auth:
            hint = "ATLASSIAN_EMAIL/ATLASSIAN_API_TOKEN are not set for this process"
        else:
            hint = (
                "verify ATLASSIAN_EMAIL/ATLASSIAN_API_TOKEN are correct, the API token "
                "has Confluence read access/scope, and this account can view the page"
            )
        raise RuntimeError(
            f"Confluence returned {resp.status_code} fetching {url} - {hint}. "
            "Note: env vars set after this app (or its terminal) was already running "
            "won't take effect until you close and restart it."
        )
    resp.raise_for_status()
    return resp.json()["body"]["view"]["value"]


def _sf_auth_error(message: str) -> RuntimeError:
    return RuntimeError(
        "Salesforce CLI has no authenticated default org.\n"
        f"Details: {message}\n"
        "Authenticate first: sf org login web --alias myorg --set-default"
    )


def _sf_is_auth_error(message: str) -> bool:
    return bool(re.search(
        r"no.*(default org|authorization|access token)|NamedOrgNotFound|NotAuthenticated|expired",
        message, re.I,
    ))


def _sf_lookup_sobject_by_prefix(sf_bin: str, record_id: str) -> str:
    """Resolve a bare record ID's object type via its 3-char key prefix."""
    import subprocess
    lookup = subprocess.run(
        [sf_bin, "data", "query", "--query",
         f"SELECT QualifiedApiName FROM EntityDefinition WHERE KeyPrefix = '{record_id[:3]}'",
         "--use-tooling-api", "--json"],
        capture_output=True, text=True, timeout=30,
    )
    lookup_data = json.loads(lookup.stdout or "{}")
    records = lookup_data.get("result", {}).get("records", [])
    if lookup.returncode != 0 or not records:
        raise ValueError(
            f"Could not determine the Salesforce object type for ID '{record_id}'. "
            "Use a full Lightning record URL "
            "(https://<domain>.lightning.force.com/lightning/r/<Object>/<Id>/view) instead."
        )
    return records[0]["QualifiedApiName"]


def _sf_fetch_record(sf_bin: str, sobject: str, record_id: str) -> tuple:
    """Fetch a single record's fields via `sf data get record`, formatted as Markdown."""
    import subprocess
    proc = subprocess.run(
        [sf_bin, "data", "get", "record", "--sobject", sobject,
         "--record-id", record_id, "--json"],
        capture_output=True, text=True, timeout=30,
    )
    data = json.loads(proc.stdout or "{}")
    if proc.returncode != 0 or data.get("status") != 0:
        message = data.get("message") or proc.stderr.strip() or "Unknown error"
        if _sf_is_auth_error(message):
            raise _sf_auth_error(message)
        raise RuntimeError(f"Salesforce CLI error fetching {sobject} {record_id}: {message}")

    fields = data.get("result", {})
    lines = [f"# Salesforce {sobject} - {record_id}", ""]
    for key, value in fields.items():
        if key == "attributes" or isinstance(value, (dict, list)) or value in (None, ""):
            continue
        lines.append(f"- **{key}**: {value}")
    return "\n".join(lines), f"{sobject}_{record_id}"


def _sf_md_cell(value) -> str:
    text = "" if value is None else str(value)
    return text.replace("\n", " ").replace("\r", " ").replace("|", "\\|")


def _sf_fetch_list_view(sf_bin: str, sobject: str, filter_name: str) -> tuple:
    """Fetch a saved list view's rows via the REST API, formatted as a Markdown table."""
    import subprocess
    # Salesforce DeveloperName/Id values are always alphanumeric + underscore; reject
    # anything else outright rather than trying to escape it into the SOQL literal.
    if not re.fullmatch(r"[A-Za-z0-9_]{1,80}", filter_name):
        raise ValueError(
            f"Invalid Salesforce list view filter name: '{filter_name}'. Expected only "
            "letters, numbers, and underscores."
        )
    # Salesforce rejects the whole query if `Id = '...'` doesn't look like an ID literal,
    # so only add that clause when filter_name is actually ID-shaped.
    if re.fullmatch(r"[a-zA-Z0-9]{15}|[a-zA-Z0-9]{18}", filter_name):
        where_clause = f"(DeveloperName = '{filter_name}' OR Id = '{filter_name}')"
    else:
        where_clause = f"DeveloperName = '{filter_name}'"
    lookup = subprocess.run(
        [sf_bin, "data", "query", "--query",
         f"SELECT Id, Name FROM ListView WHERE SobjectType = '{sobject}' AND {where_clause}",
         "--json"],
        capture_output=True, text=True, timeout=30,
    )
    lookup_data = json.loads(lookup.stdout or "{}")
    lv_records = lookup_data.get("result", {}).get("records", [])
    if lookup.returncode != 0 or not lv_records:
        raise ValueError(
            f"Could not find a saved list view '{filter_name}' on {sobject}. "
            "Standard/system list views (e.g. \"Recent\", \"All Open Cases\") aren't "
            "queryable this way - open a custom list view instead, or use a single "
            "record URL."
        )
    list_view_id = lv_records[0]["Id"]
    list_view_label = lv_records[0].get("Name", filter_name)

    org_info = subprocess.run(
        [sf_bin, "org", "display", "--json"], capture_output=True, text=True, timeout=30,
    )
    try:
        api_version = json.loads(org_info.stdout or "{}")["result"]["apiVersion"]
    except (json.JSONDecodeError, KeyError):
        api_version = "60.0"

    proc = subprocess.run(
        [sf_bin, "api", "request", "rest",
         f"services/data/v{api_version}/sobjects/{sobject}/listviews/{list_view_id}/results?pageSize=2000"],
        capture_output=True, text=True, timeout=60,
    )
    try:
        results = json.loads(proc.stdout or "null")
    except json.JSONDecodeError:
        results = None

    if proc.returncode != 0:
        message = proc.stderr.strip() or "Unknown error"
        if isinstance(results, list) and results:
            message = results[0].get("message", message)
        elif isinstance(results, dict):
            message = results.get("message", message)
        if _sf_is_auth_error(message):
            raise _sf_auth_error(message)
        raise RuntimeError(f"Salesforce CLI error fetching list view '{filter_name}': {message}")
    if results is None:
        raise RuntimeError(
            f"Salesforce CLI returned an unparseable response for list view '{filter_name}'."
        )

    columns = [_sf_md_cell(c.get("label", "")) for c in results.get("columns", [])]
    records = results.get("records", [])
    total = results.get("size", len(records))

    lines = [f"# Salesforce {sobject} list view - {list_view_label}", ""]
    lines.append(
        f"*Showing {len(records)} of {total} records (truncated).*\n"
        if total > len(records) else f"*{len(records)} record(s).*\n"
    )
    lines.append("| " + " | ".join(columns) + " |")
    lines.append("|" + "|".join(["---"] * len(columns)) + "|")
    for rec in records:
        values = [_sf_md_cell(col.get("value")) for col in rec.get("columns", [])]
        lines.append("| " + " | ".join(values) + " |")

    return "\n".join(lines), f"{sobject}_listview_{filter_name}"


def _fetch_salesforce(url: str) -> tuple:
    """Fetch a Salesforce record or list view via the authenticated `sf` CLI, as Markdown.

    A plain HTTP GET only returns Salesforce's login page (no session cookie),
    so content goes through the Salesforce CLI's already-authenticated default
    org instead of the requests/trafilatura path used for other URLs.
    """
    import shutil
    from urllib.parse import urlparse, parse_qs

    sf_bin = shutil.which("sf")
    if not sf_bin:
        raise RuntimeError(
            "Salesforce CLI ('sf') not found on PATH - Salesforce URLs need an "
            "authenticated CLI session because a plain HTTP request only returns the "
            "login page.\n"
            "Install: npm install --global @salesforce/cli\n"
            "Then authenticate once: sf org login web --alias myorg --set-default"
        )

    parsed = urlparse(url)
    path = parsed.path

    m = re.search(r"/lightning/r/([A-Za-z0-9_]+)/([a-zA-Z0-9]{15,18})(?:/|$)", path)
    if m:
        sobject, record_id = m.group(1), m.group(2)
        return _sf_fetch_record(sf_bin, sobject, record_id)

    m = re.fullmatch(r"/?([a-zA-Z0-9]{15}|[a-zA-Z0-9]{18})/?", path)
    if m:
        record_id = m.group(1)
        sobject = _sf_lookup_sobject_by_prefix(sf_bin, record_id)
        return _sf_fetch_record(sf_bin, sobject, record_id)

    m = re.search(r"/lightning/o/([A-Za-z0-9_]+)/list", path)
    if m:
        sobject = m.group(1)
        filter_name = parse_qs(parsed.query).get("filterName", [None])[0]
        if not filter_name:
            raise ValueError(
                "This Salesforce list view URL doesn't specify a filterName, so SlimDocs "
                "can't tell which saved list view to fetch. Open a specific list view "
                "(URL contains '?filterName=...') or paste a single record URL instead."
            )
        return _sf_fetch_list_view(sf_bin, sobject, filter_name)

    raise ValueError(f"Could not find a Salesforce record ID or list view in URL: {url}")


_JS_REQUIRED_RE = re.compile(
    r"you (need|must) (to )?enable javascript|javascript is required|"
    r"please enable javascript|this (app|page|site) requires javascript|"
    r"enable javascript to (run|view|continue)",
    re.IGNORECASE,
)
_REDIRECT_STUB_RE = re.compile(
    r"if you are not redirected|having trouble accessing|click here to continue",
    re.IGNORECASE,
)
_META_REFRESH_RE = re.compile(r'<meta[^>]+http-equiv=["\']?refresh["\']?', re.IGNORECASE)


def _fetch_url(url: str, fmt: str) -> tuple:
    """Fetch a URL and return (extracted_text, source_name, raw_html_bytes)."""
    from urllib.parse import urlparse, unquote
    try:
        import requests
    except ImportError:
        raise RuntimeError("requests not installed - run: pip install requests")
    try:
        import trafilatura
    except ImportError:
        raise RuntimeError("trafilatura not installed - run: pip install trafilatura")

    parsed = urlparse(url)
    path_stem = Path(unquote(parsed.path)).stem or Path(unquote(parsed.path)).parent.name or "page"
    source_name = re.sub(r"[^\w\-]", "_", f"{parsed.netloc}_{path_stem}")[:80].strip("_")

    if "salesforce.com" in parsed.netloc or "force.com" in parsed.netloc:
        content, sf_source = _fetch_salesforce(url)
        return content, sf_source, len(content.encode("utf-8"))

    if "atlassian.net" in parsed.netloc and "/pages/" in parsed.path:
        html = _fetch_confluence(url)
    else:
        resp = requests.get(url, timeout=20, headers={
            "User-Agent": "Mozilla/5.0 (compatible; TokenReducer/1.0)"
        })
        resp.raise_for_status()
        html = resp.text

    raw_html_bytes = len(html.encode("utf-8"))
    trafilatura_fmt = "txt" if fmt == "txt" else "markdown"
    content = trafilatura.extract(
        html,
        output_format=trafilatura_fmt,
        include_links=(fmt != "txt"),
        include_images=False,
        favor_recall=True,
    )
    if not content or not content.strip():
        raise ValueError("No content could be extracted from this URL")
    if len(content.strip()) < 300:
        if _JS_REQUIRED_RE.search(content):
            raise ValueError(
                "This page appears to be a JavaScript-rendered single-page app - only a "
                "'JavaScript is required' placeholder was found. SlimDocs fetches static "
                "HTML only and can't execute JavaScript, so this page's real content can't "
                "be extracted this way."
            )
        if _META_REFRESH_RE.search(html) or _REDIRECT_STUB_RE.search(content):
            raise ValueError(
                "This page is an automatic-redirect interstitial, not real content - only "
                "redirect/placeholder text was found (e.g. 'click here if you are not "
                "redirected'). SlimDocs fetches static HTML only and doesn't follow "
                "client-side/JavaScript redirects, so it can't reach the actual "
                "destination page this way."
            )
    return content, source_name, raw_html_bytes


# ── Text processing ───────────────────────────────────────────────────────────

def chunk_text(text: str, max_chars: int = 2000, overlap: int = 200) -> list:
    chunks, start = [], 0
    length = len(text)
    while start < length:
        chunks.append(text[start:start + max_chars])
        start += max_chars - overlap
    return chunks


def _make_duckdb_db(output_dir: Path) -> tuple:
    """Create a timestamped .duckdb file, initialise the documents table, return (db_path, conn)."""
    try:
        import duckdb
    except ImportError:
        raise RuntimeError("duckdb not installed - run: pip install duckdb")
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    db_path = output_dir / f"slimDocs_{ts}.duckdb"
    conn = duckdb.connect(str(db_path))
    conn.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            doc_id        INTEGER PRIMARY KEY,
            source        VARCHAR NOT NULL,
            name          VARCHAR,
            extension     VARCHAR,
            content       TEXT,
            tokens_before INTEGER,
            tokens_after  INTEGER,
            tokens_saved  INTEGER,
            size_kb       DOUBLE,
            extracted_at  TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    return db_path, conn


def process_files_duckdb(
    files: list,
    output_dir: Path,
    root: "Path | None",
    progress_bar,
    status_text,
    content_map: dict,
) -> dict:
    _pdf_img_cache.clear()
    successes, errors = [], []
    total = len(files)
    db_path, conn = _make_duckdb_db(output_dir)
    doc_id = 1

    try:
        for i, file_path in enumerate(files):
            progress_bar.progress((i + 1) / total)
            status_text.text(f"Processing {i + 1}/{total}: {file_path.name}")
            try:
                content = extract_from_file(file_path)
                if not content or not content.strip():
                    errors.append({"File": str(file_path), "Error": "No content extracted"})
                    continue

                raw_size  = file_path.stat().st_size
                out_chars = len(content)
                raw_tok   = raw_size  // 4
                out_tok   = out_chars // 4
                saved     = max(0, raw_tok - out_tok)
                pct_red   = (1 - out_chars / raw_size) * 100 if raw_size else 0.0

                conn.execute(
                    """
                    INSERT INTO documents (doc_id, source, name, extension, content,
                                          tokens_before, tokens_after, tokens_saved, size_kb)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    [doc_id, str(file_path), file_path.name,
                     file_path.suffix.lower(), content,
                     raw_tok, out_tok, saved, round(raw_size / 1024, 1)],
                )
                doc_id += 1
                _pdf_img_cache.pop(str(file_path), None)

                content_map[str(file_path)] = {"raw": content, "formatted": content, "fmt": "md"}
                successes.append({
                    "File":          file_path.name,
                    "Output":        "DuckDB",
                    "Size (KB)":     round(raw_size / 1024, 1),
                    "Tokens Before": raw_tok,
                    "Tokens After":  out_tok,
                    "Saved":         saved,
                    "Reduction %":   round(pct_red, 1),
                    "_path":         str(file_path),
                    "_ext":          file_path.suffix.lower(),
                })
            except Exception as e:
                errors.append({"File": str(file_path), "Error": str(e)})
    finally:
        conn.close()

    run_ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    return {
        "successes":   successes,
        "errors":      errors,
        "output_path": str(db_path),
        "zip_output":  False,
        "is_duckdb":   True,
        "run_ts":      run_ts,
        "logs":        _build_logs(successes, errors, run_ts),
    }


def process_urls_duckdb(
    urls: list,
    output_dir: Path,
    progress_bar,
    status_text,
    content_map: dict,
) -> dict:
    successes, errors = [], []
    total = len(urls)
    used_names: set = set()
    db_path, conn = _make_duckdb_db(output_dir)
    doc_id = 1

    try:
        for i, url in enumerate(urls):
            progress_bar.progress((i + 1) / total)
            short = url[:65] + "…" if len(url) > 65 else url
            status_text.text(f"Fetching {i + 1}/{total}: {short}")
            try:
                content, source_name, raw_bytes = _fetch_url(url, "md")

                raw_tok = raw_bytes  // 4
                out_tok = len(content) // 4
                saved   = max(0, raw_tok - out_tok)
                pct_red = (1 - out_tok / max(raw_tok, 1)) * 100

                name, counter = source_name, 1
                while name in used_names:
                    name = f"{source_name}_{counter}"
                    counter += 1
                used_names.add(name)

                conn.execute(
                    """
                    INSERT INTO documents (doc_id, source, name, extension, content,
                                          tokens_before, tokens_after, tokens_saved, size_kb)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    [doc_id, url, name, ".url", content,
                     raw_tok, out_tok, saved, round(raw_bytes / 1024, 1)],
                )
                doc_id += 1

                content_map[url] = {"raw": content, "formatted": content, "fmt": "md"}
                successes.append({
                    "File":          name,
                    "URL":           url,
                    "Output":        "DuckDB",
                    "Size (KB)":     round(raw_bytes / 1024, 1),
                    "Tokens Before": raw_tok,
                    "Tokens After":  out_tok,
                    "Saved":         saved,
                    "Reduction %":   round(pct_red, 1),
                    "_path":         url,
                    "_ext":          ".url",
                })
            except Exception as e:
                errors.append({"File": url, "Error": str(e)})
    finally:
        conn.close()

    run_ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    return {
        "successes":   successes,
        "errors":      errors,
        "output_path": str(db_path),
        "zip_output":  False,
        "is_duckdb":   True,
        "run_ts":      run_ts,
        "logs":        _build_logs(successes, errors, run_ts),
    }


def _format_output(content: str, fmt: str, stem: str, original_source) -> tuple:
    """Return (text_to_write, file_extension)."""
    header = f"Source: {original_source}\n\n"
    if fmt == "md":
        return header + content, "md"
    elif fmt == "txt":
        return header + content, "txt"
    else:
        chunks = chunk_text(content)
        payload = [
            {
                "id": f"{stem}_chunk_{i}",
                "content": chunk,
                "metadata": {"source": str(original_source)},
            }
            for i, chunk in enumerate(chunks)
        ]
        return json.dumps(payload, indent=2, ensure_ascii=False), "json"


# ── File discovery ────────────────────────────────────────────────────────────

def discover_files(root: Path, recursive: bool = True) -> list:
    pattern = "**/*" if recursive else "*"
    return sorted(
        p for p in root.glob(pattern)
        if p.is_file() and p.suffix.lower() in _SUPPORTED_EXTENSIONS
    )


def _discover_skipped_files(root: Path, recursive: bool = True) -> list:
    """Files in the folder excluded from discover_files() due to an unsupported extension."""
    pattern = "**/*" if recursive else "*"
    return sorted(
        p for p in root.glob(pattern)
        if p.is_file() and p.suffix.lower() not in _SUPPORTED_EXTENSIONS
    )


# ── Native dialog helpers ─────────────────────────────────────────────────────

def _pick_folder() -> str | None:
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk(); root.withdraw(); root.wm_attributes("-topmost", 1)
        folder = filedialog.askdirectory(title="Select folder")
        root.destroy()
        return folder or None
    except Exception:
        return None


def _pick_file() -> str | None:
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk(); root.withdraw(); root.wm_attributes("-topmost", 1)
        ext_list = " ".join(f"*{e}" for e in sorted(_SUPPORTED_EXTENSIONS))
        path = filedialog.askopenfilename(
            title="Select file",
            filetypes=[("Supported files", ext_list), ("All files", "*.*")],
        )
        root.destroy()
        return path or None
    except Exception:
        return None


# ── Keyboard fix ──────────────────────────────────────────────────────────────

def _inject_keyboard_fix():
    st.html("""
        <script>
        (function () {
            const EDIT_KEYS = new Set(['c', 'x', 'v', 'a', 'z', 'y']);
            function blockInTextArea(e) {
                if (!(e.ctrlKey || e.metaKey)) return;
                if (!EDIT_KEYS.has(e.key.toLowerCase())) return;
                const el = e.target || document.activeElement;
                if (!el) return;
                const tag = (el.tagName || '').toLowerCase();
                if (tag === 'textarea' || tag === 'input' || el.isContentEditable)
                    e.stopImmediatePropagation();
            }
            document.addEventListener('keydown', blockInTextArea, true);
        })();
        </script>
    """)


# ── Core processing ───────────────────────────────────────────────────────────

def _make_bundle(output_dir: Path, zip_output: bool) -> tuple:
    """Return (bundle_name, zip_path_or_None, out_folder_or_None, zf_or_None)."""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    bundle_name = f"extracted_{ts}"
    if zip_output:
        zip_path = output_dir / f"{bundle_name}.zip"
        zf = zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED)
        return bundle_name, zip_path, None, zf
    else:
        out_folder = output_dir / bundle_name
        out_folder.mkdir(parents=True, exist_ok=True)
        return bundle_name, None, out_folder, None


def _build_logs(successes: list, errors: list, run_ts: str) -> list:
    logs = []
    for r in successes:
        logs.append({
            "Time": run_ts,
            "Type": "✅ Success",
            "File": r["File"],
            "Message": (
                f"{r['Tokens Before']:,} tokens → {r['Tokens After']:,} tokens  "
                f"({r['Reduction %']:.1f}% reduction, saved {r['Saved']:,})"
            ),
        })
    for e in errors:
        logs.append({
            "Time": run_ts,
            "Type": "❌ Error",
            "File": Path(e["File"]).name if not e["File"].startswith("http") else e["File"][:60],
            "Message": e["Error"],
        })
    return logs


def _record_error_only_run(errors: list, output_path: str, zip_output: bool) -> dict:
    """Record a run that produced no successes, so validation failures (unsupported
    file type, nothing left to process, etc.) still show up in Logs/Statistics
    instead of only flashing an inline error in File Processing.
    """
    run_ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    result = {
        "successes":   [],
        "errors":      errors,
        "output_path": output_path,
        "zip_output":  zip_output,
        "run_ts":      run_ts,
        "logs":        _build_logs([], errors, run_ts),
    }
    st.session_state.results = result
    st.session_state.logs    = result["logs"] + st.session_state.logs
    st.session_state.session_history[run_ts] = result
    st.session_state.selected_session_ts = None
    st.session_state.session_totals["runs"]   += 1
    st.session_state.session_totals["errors"] += len(errors)
    return result


def process_files(
    files: list,
    fmt: str,
    output_dir: Path,
    zip_output: bool,
    root: Path | None,
    progress_bar,
    status_text,
    content_map: dict,
) -> dict:
    _pdf_img_cache.clear()
    successes, errors = [], []
    total = len(files)
    _, zip_path, out_folder, zf = _make_bundle(output_dir, zip_output)

    try:
        for i, file_path in enumerate(files):
            progress_bar.progress((i + 1) / total)
            status_text.text(f"Processing {i + 1}/{total}: {file_path.name}")
            try:
                content = extract_from_file(file_path)
                if not content or not content.strip():
                    errors.append({"File": str(file_path), "Error": "No content extracted"})
                    continue

                raw_size  = file_path.stat().st_size
                out_chars = len(content)
                raw_tok   = raw_size  // 4
                out_tok   = out_chars // 4
                saved     = max(0, raw_tok - out_tok)
                pct_red   = (1 - out_chars / raw_size) * 100 if raw_size else 0.0

                text_out, ext = _format_output(content, fmt, file_path.stem, file_path)
                out_name = (
                    str(file_path.relative_to(root).with_suffix(f".{ext}"))
                    if root else f"{file_path.stem}.{ext}"
                )

                if zf:
                    zf.writestr(out_name, text_out)
                else:
                    out_p = out_folder / out_name
                    out_p.parent.mkdir(parents=True, exist_ok=True)
                    out_p.write_text(text_out, encoding="utf-8")

                # Write companion images extracted from PDF (populated by _extract_pdf_images_ocr)
                out_parent = Path(out_name).parent  # "" for root, "subdir" for nested
                for rel_img, img_bytes in _pdf_img_cache.pop(str(file_path), []):
                    full_img_rel = str(out_parent / rel_img) if str(out_parent) != "." else rel_img
                    if zf:
                        zf.writestr(full_img_rel, img_bytes)
                    else:
                        img_p = out_folder / full_img_rel
                        img_p.parent.mkdir(parents=True, exist_ok=True)
                        img_p.write_bytes(img_bytes)

                content_map[str(file_path)] = {"raw": content, "formatted": text_out, "fmt": ext}
                successes.append({
                    "File":          file_path.name,
                    "Output":        out_name,
                    "Size (KB)":     round(raw_size / 1024, 1),
                    "Tokens Before": raw_tok,
                    "Tokens After":  out_tok,
                    "Saved":         saved,
                    "Reduction %":   round(pct_red, 1),
                    "_path":         str(file_path),
                    "_ext":          file_path.suffix.lower(),
                })
            except Exception as e:
                errors.append({"File": str(file_path), "Error": str(e)})
    finally:
        if zf:
            zf.close()

    run_ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    output_path = str(zip_path) if zip_output else str(out_folder)
    return {
        "successes":   successes,
        "errors":      errors,
        "output_path": output_path,
        "zip_output":  zip_output,
        "run_ts":      run_ts,
        "logs":        _build_logs(successes, errors, run_ts),
    }


def process_urls(
    urls: list,
    fmt: str,
    output_dir: Path,
    zip_output: bool,
    progress_bar,
    status_text,
    content_map: dict,
) -> dict:
    successes, errors = [], []
    total = len(urls)
    used_names: set = set()
    _, zip_path, out_folder, zf = _make_bundle(output_dir, zip_output)

    try:
        for i, url in enumerate(urls):
            progress_bar.progress((i + 1) / total)
            short = url[:65] + "…" if len(url) > 65 else url
            status_text.text(f"Fetching {i + 1}/{total}: {short}")
            try:
                content, source_name, raw_bytes = _fetch_url(url, fmt)

                raw_tok = raw_bytes  // 4
                out_tok = len(content) // 4
                saved   = max(0, raw_tok - out_tok)
                pct_red = (1 - out_tok / max(raw_tok, 1)) * 100

                text_out, ext = _format_output(content, fmt, source_name, url)
                base = f"{source_name}.{ext}"
                out_name, counter = base, 1
                while out_name in used_names:
                    out_name = f"{source_name}_{counter}.{ext}"
                    counter += 1
                used_names.add(out_name)

                if zf:
                    zf.writestr(out_name, text_out)
                else:
                    out_p = out_folder / out_name
                    out_p.write_text(text_out, encoding="utf-8")

                content_map[url] = {"raw": content, "formatted": text_out, "fmt": ext}
                successes.append({
                    "File":          source_name,
                    "URL":           url,
                    "Output":        out_name,
                    "Size (KB)":     round(raw_bytes / 1024, 1),
                    "Tokens Before": raw_tok,
                    "Tokens After":  out_tok,
                    "Saved":         saved,
                    "Reduction %":   round(pct_red, 1),
                    "_path":         url,
                    "_ext":          f".{ext}",
                })
            except Exception as e:
                errors.append({"File": url, "Error": str(e)})
    finally:
        if zf:
            zf.close()

    run_ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    output_path = str(zip_path) if zip_output else str(out_folder)
    return {
        "successes":   successes,
        "errors":      errors,
        "output_path": output_path,
        "zip_output":  zip_output,
        "run_ts":      run_ts,
        "logs":        _build_logs(successes, errors, run_ts),
    }


# ── DOCX formatted preview ────────────────────────────────────────────────────

def _render_docx_preview(path: Path):
    """Render a DOCX file as structured Streamlit content (headings, runs, tables)."""
    try:
        from docx import Document
        from docx.oxml.ns import qn as _qn
        from docx.text.paragraph import Paragraph as _Para
        from docx.table import Table as _Table
    except ImportError:
        st.info("Install python-docx (`pip install python-docx`) for DOCX preview.")
        return

    doc = Document(path)

    _HEADING_PREFIXES = {
        "title":      "#",
        "subtitle":   "##",
        "heading 1":  "#",
        "heading 2":  "##",
        "heading 3":  "###",
        "heading 4":  "####",
        "heading 5":  "#####",
    }

    def _fmt_runs(para) -> str:
        parts = []
        for run in para.runs:
            t = run.text
            if not t:
                continue
            if run.bold and run.italic:
                t = f"***{t}***"
            elif run.bold:
                t = f"**{t}**"
            elif run.italic:
                t = f"*{t}*"
            parts.append(t)
        return "".join(parts) or para.text

    def _render_para(para):
        text = para.text.strip()
        if not text:
            return
        style_key = para.style.name.lower()
        prefix = next(
            (v for k, v in _HEADING_PREFIXES.items() if style_key.startswith(k)),
            None,
        )
        if prefix:
            st.markdown(f"{prefix} {text}")
        else:
            formatted = _fmt_runs(para).strip()
            if formatted:
                st.markdown(formatted)

    def _render_table(tbl):
        rows = [
            [cell.text.strip() for cell in row.cells]
            for row in tbl.rows
        ]
        if not rows:
            return
        st.divider()
        if len(rows) > 1:
            seen: dict = {}
            headers = []
            for h in rows[0]:
                key = h or "Col"
                if key in seen:
                    seen[key] += 1
                    key = f"{key}_{seen[key]}"
                else:
                    seen[key] = 0
                headers.append(key)
            df = pd.DataFrame(rows[1:], columns=headers)
        else:
            df = pd.DataFrame(rows)
        st.dataframe(df, use_container_width=True, hide_index=True)
        st.divider()

    # Walk body children in document order so paragraphs and tables appear correctly
    for child in doc.element.body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            _render_para(_Para(child, doc))
        elif tag == "tbl":
            _render_table(_Table(child, doc))


# ── PPTX formatted preview ────────────────────────────────────────────────────

def _render_pptx_preview(path: Path):
    """Render a PPTX file slide-by-slide: titles, bullets, tables, images, notes."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[attr-defined]
    except ImportError:
        st.info("Install python-pptx (`pip install python-pptx`) for PPTX preview.")
        return

    prs        = Presentation(path)
    all_slides = list(prs.slides)   # materialise - prs.slides doesn't support slicing
    total      = len(all_slides)
    _MAX_SLIDES = 50
    shown = min(total, _MAX_SLIDES)
    st.caption(
        f"PowerPoint · {total} slide(s)"
        + (f" · showing first {shown}" if total > _MAX_SLIDES else "")
    )

    for slide_num, slide in enumerate(all_slides[:shown], 1):
        # Build the expander label from the slide title (if available)
        title_shape = slide.shapes.title
        title_text  = title_shape.text.strip() if title_shape and title_shape.text.strip() else ""
        label = f"Slide {slide_num}" + (f" - {title_text}" if title_text else "")

        with st.expander(label, expanded=(slide_num == 1)):
            if title_text:
                st.markdown(f"### {title_text}")

            for shape in slide.shapes:
                # ── Skip the title shape (already rendered above) ─────────
                if shape == title_shape:
                    continue

                # ── Inline picture ────────────────────────────────────────
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        blob = shape.image.blob
                        st.image(blob, use_container_width=False)
                        continue
                except Exception:
                    pass

                # ── Table ─────────────────────────────────────────────────
                if shape.has_table:
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if rows:
                        if len(rows) > 1:
                            # Deduplicate header cells (merged/empty cells cause Arrow errors)
                            seen: dict = {}
                            headers = []
                            for h in rows[0]:
                                key = h or "Col"
                                if key in seen:
                                    seen[key] += 1
                                    key = f"{key}_{seen[key]}"
                                else:
                                    seen[key] = 0
                                headers.append(key)
                            df = pd.DataFrame(rows[1:], columns=headers)
                        else:
                            df = pd.DataFrame(rows)
                        st.dataframe(df, use_container_width=True, hide_index=True)
                    continue

                # ── Text frame (bullets / body text) ─────────────────────
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        level = para.level or 0
                        if level == 0:
                            st.markdown(text)
                        else:
                            st.markdown("&nbsp;" * (level * 4) + f"• {text}",
                                        unsafe_allow_html=True)

            # ── Speaker notes ─────────────────────────────────────────────
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        st.caption(f"📝 Notes: {notes}")
            except Exception:
                pass

    if total > _MAX_SLIDES:
        st.info(f"… {total - _MAX_SLIDES} more slides not shown.")


# ── Original-file renderer ────────────────────────────────────────────────────

def _render_original_file(path_str: str):
    """Render a best-effort in-app view of the source file."""
    # URLs - nothing to render locally
    if path_str.startswith("http://") or path_str.startswith("https://"):
        st.info(f"Source is a URL: {path_str}")
        return

    path = Path(path_str)
    if not path.exists():
        st.warning("Original file is no longer accessible on disk.")
        return

    ext = path.suffix.lower()

    # ── Raster images ─────────────────────────────────────────────────────
    if ext in _EXT_CATEGORIES["Images"]:
        st.image(str(path), use_container_width=True)
        return

    # ── PDF - render each page as an image via PyMuPDF ───────────────────
    if ext == ".pdf":
        try:
            import fitz
        except ImportError:
            st.info(
                "Install **PyMuPDF** to enable inline PDF preview:  \n"
                "`pip install PyMuPDF`"
            )
            return
        _MAX_PDF_PAGES = 20
        doc = fitz.open(str(path))
        total = len(doc)
        show  = min(total, _MAX_PDF_PAGES)
        st.caption(
            f"PDF · {total} page(s)"
            + (f" · showing first {show}" if total > _MAX_PDF_PAGES else "")
        )
        for pg_num in range(show):
            page = doc[pg_num]
            pix  = page.get_pixmap(dpi=130)
            st.image(pix.tobytes("png"), caption=f"Page {pg_num + 1}",
                     use_container_width=True)
        if total > _MAX_PDF_PAGES:
            st.info(f"… {total - _MAX_PDF_PAGES} more pages not shown.")
        return

    # ── Plain-text families (code, config, log, md, txt, html, csv …) ─────
    _TEXT_PREVIEW_EXTS = (
        _PLAINTEXT_EXTS
        | {".md", ".txt", ".html", ".htm", ".csv", ".tsv"}
        | _EXT_CATEGORIES["Config/Data"]
        | _EXT_CATEGORIES["Code"]
        | _EXT_CATEGORIES["Logs"]
    )
    if ext in _TEXT_PREVIEW_EXTS:
        try:
            raw = _decode_bytes(path.read_bytes())
        except Exception as e:
            st.error(f"Could not read file: {e}")
            return
        trunc = len(raw) > _PREVIEW_LIMIT
        if trunc:
            st.caption(f"Showing first {_PREVIEW_LIMIT:,} of {len(raw):,} characters.")
        _EXT_LANG = {
            ".py": "python", ".pyw": "python",
            ".js": "javascript", ".mjs": "javascript", ".cjs": "javascript",
            ".ts": "typescript", ".tsx": "typescript", ".jsx": "javascript",
            ".java": "java", ".kt": "kotlin", ".go": "go",
            ".rs": "rust", ".cs": "csharp", ".vb": "vb",
            ".cpp": "cpp", ".cxx": "cpp", ".cc": "cpp", ".c": "c", ".h": "c",
            ".hpp": "cpp", ".hxx": "cpp",
            ".rb": "ruby", ".php": "php", ".swift": "swift",
            ".sh": "bash", ".bash": "bash", ".zsh": "bash", ".fish": "bash",
            ".bat": "batch", ".cmd": "batch",
            ".ps1": "powershell", ".psm1": "powershell",
            ".sql": "sql", ".ddl": "sql", ".dml": "sql",
            ".html": "html", ".htm": "html",
            ".xml": "xml", ".yaml": "yaml", ".yml": "yaml",
            ".json": "json", ".jsonc": "json", ".json5": "json",
            ".toml": "toml", ".md": "markdown",
            ".css": "css", ".lua": "lua",
        }
        lang = _EXT_LANG.get(ext, "text")
        st.code(raw[:_PREVIEW_LIMIT] if trunc else raw, language=lang, wrap_lines=True)
        return

    # ── EVTX - show parsed event records inline ───────────────────────────
    if ext == ".evtx":
        entry = st.session_state.content_map.get(path_str, {})
        raw = entry.get("raw", "")
        if raw:
            trunc = len(raw) > _PREVIEW_LIMIT
            if trunc:
                st.caption(f"Showing first {_PREVIEW_LIMIT:,} of {len(raw):,} characters.")
            st.code(raw[:_PREVIEW_LIMIT] if trunc else raw, language="text", wrap_lines=True)
        else:
            st.info("Run extraction to view parsed event log records here.")
        return

    # ── Archives - file listing + extracted text ──────────────────────────
    if ext in _EXT_CATEGORIES["Archives"]:
        if ext == ".zip":
            try:
                import zipfile, pandas as _pd
                with zipfile.ZipFile(path, "r") as _zf:
                    infos = _zf.infolist()
                st.caption(f"ZIP archive · {len(infos):,} item(s)")
                rows = [
                    {
                        "Name":              i.filename,
                        "Type":              "Dir" if i.is_dir() else "File",
                        "Size (KB)":         round(i.file_size   / 1024, 1),
                        "Compressed (KB)":   round(i.compress_size / 1024, 1),
                    }
                    for i in infos
                ]
                st.dataframe(_pd.DataFrame(rows), use_container_width=True, hide_index=True)
            except Exception as _e:
                st.warning(f"Could not list ZIP contents: {_e}")
        entry = st.session_state.content_map.get(path_str, {})
        raw = entry.get("raw", "")
        if raw:
            st.divider()
            st.caption("Extracted text content from archive:")
            trunc = len(raw) > _PREVIEW_LIMIT
            if trunc:
                st.caption(f"Showing first {_PREVIEW_LIMIT:,} of {len(raw):,} characters.")
            st.code(raw[:_PREVIEW_LIMIT] if trunc else raw, language="text", wrap_lines=True)
        elif ext != ".zip":
            st.info("Run extraction to view archive contents here.")
        return

    # ── Videos ────────────────────────────────────────────────────────────
    if ext in _EXT_CATEGORIES["Videos"]:
        stat = path.stat()
        st.caption(
            f"{path.name}  ·  {ext.upper().lstrip('.')}  ·  "
            f"{stat.st_size / 1_048_576:.2f} MB"
        )
        try:
            st.video(str(path))
        except Exception:
            st.info(
                "Video preview unavailable.  \n"
                "The Extracted text tab shows file metadata only - "
                "audio/video transcription is not supported."
            )
        return

    # ── DOCX - formatted render ───────────────────────────────────────────
    if ext in {".docx"}:
        _render_docx_preview(path)
        return

    if ext == ".doc":
        st.info(
            "`.doc` is the legacy Word binary format - python-docx cannot open it.  \n"
            "Convert to `.docx` in Word, or see the **Extracted text** tab."
        )
        return

    # ── PPTX - slide-by-slide render ─────────────────────────────────────
    if ext == ".pptx":
        _render_pptx_preview(path)
        return

    if ext == ".ppt":
        st.info(
            "`.ppt` is the legacy PowerPoint binary format - python-pptx cannot open it.  \n"
            "Convert to `.pptx` in PowerPoint, or see the **Extracted text** tab."
        )
        return

    # ── Binary documents (XLSX, ODT, ODS …) ──────────────────────────────
    st.info(
        f"In-app preview is not available for `{ext}` files. "
        "See the **Extracted text** tab for the parsed content."
    )


# ── Content preview panel ─────────────────────────────────────────────────────

def _show_preview(row: dict):
    path_str = row["_path"]
    entry = st.session_state.content_map.get(path_str, {})
    if not entry:
        st.warning("Content not available for this file.")
        return

    raw      = entry.get("raw", "")
    formatted = entry.get("formatted", "")
    fmt_ext  = entry.get("fmt", "md")
    lang_map = {"md": "markdown", "json": "json", "txt": "text"}
    lang     = lang_map.get(fmt_ext, "text")

    label = row.get("URL", row["File"]) if "URL" in row else row["File"]
    st.markdown(f"**Preview - {label}**")

    tab_orig, tab_raw, tab_fmt = st.tabs([
        "📂 Original file",
        f"📄 Extracted text  ({row['Tokens Before']:,} tok before)",
        f"📝 Output .{fmt_ext}  ({row['Tokens After']:,} tok after)",
    ])
    with tab_orig:
        _render_original_file(path_str)
    with tab_raw:
        trunc = len(raw) > _PREVIEW_LIMIT
        if trunc:
            st.caption(f"Showing first {_PREVIEW_LIMIT:,} of {len(raw):,} characters.")
        st.code(raw[:_PREVIEW_LIMIT] if trunc else raw, language="markdown", wrap_lines=True)
    with tab_fmt:
        trunc = len(formatted) > _PREVIEW_LIMIT
        if trunc:
            st.caption(f"Showing first {_PREVIEW_LIMIT:,} of {len(formatted):,} characters.")
        st.code(formatted[:_PREVIEW_LIMIT] if trunc else formatted, language=lang, wrap_lines=True)


# ── Help dialog ───────────────────────────────────────────────────────────────

@st.dialog("📄 SlimDocs - Help", width="large")
def _show_help():
    st.markdown("""
## How to use

| Step | Action |
|------|--------|
| **1. Input mode** | Single File, Folder, or URLs |
| **2. Format** | `md` Markdown (default) · `txt` Plain text · `json` Chunked Claude JSON · `duckdb` Queryable database |
| **3. Output dir** | Defaults to Downloads |
| **4. Zip toggle** | ON = single `.zip`  OFF = timestamped folder  *(hidden for DuckDB)* |
| **5. Extract** | Click 🚀 |

---

## Supported file types

| Category | Extensions |
|----------|-----------|
| **Documents** | `.pdf` `.docx` `.doc` `.pptx` `.ppt` `.odt` `.rtf` `.md` `.txt` |
| **Spreadsheets** | `.xlsx` `.xls` `.csv` `.tsv` `.ods` |
| **Web** | `.html` `.htm` |
| **Archives** | `.zip` `.tar` `.gz` `.bz2` `.rar` `.xz` `.7z` |
| **Images (OCR)** | `.png` `.jpg` `.jpeg` `.gif` `.bmp` `.tiff` `.webp` |
| **Logs** | `.log` `.out` `.err` `.evtx` |
| **Code** | `.py` `.js` `.ts` `.java` `.go` `.rs` `.cs` `.cpp` `.sql` `.sh` `.ps1` … |
| **Config/Data** | `.json` `.yaml` `.toml` `.xml` `.ini` `.env` `.tf` `.tfvars` … |
| **Videos** | `.mp4` `.avi` `.mov` `.mkv` … (file metadata only) |

---

## URL mode

Paste one or more URLs (one per line or comma-separated).

- **Confluence** pages are fetched via the REST API if `ATLASSIAN_EMAIL` and `ATLASSIAN_API_TOKEN` env vars are set.
- **Salesforce** record, bare-ID, and custom list-view URLs (`*.salesforce.com` / `*.lightning.force.com`) are fetched via the `sf` CLI, since a plain request only returns the login page. Requires the CLI installed (`npm install --global @salesforce/cli`) and an authenticated default org (`sf org login web -a UiPath --set-default`) - see the README for full setup steps.
- **JavaScript-rendered / redirect-only pages** (SPA shells, "click here if you are not redirected" interstitials) are detected and reported as errors instead of silently extracting placeholder text - SlimDocs fetches static HTML only.

---

## Statistics & Reports and Logs

- **Statistics & Reports** shows charts, a full results table, and an **Errors** table (even for runs that only errored) for the run currently being viewed, plus a running Session Total once you've done more than one run.
- **Logs** keeps a full history of every extraction run. Click a log row, then open **Statistics & Reports** yourself to see that run's report, with a **⬅️ Back to latest** button to return.
- Choosing the **duckdb** output format unlocks a **DuckDB Explorer** tab: keyword search with highlighted snippets, a raw SQL editor, and a table browser - queried directly from the `.duckdb` file.

---

## Optional dependencies

```
pip install pdfplumber          # PDF text/tables
pip install PyMuPDF             # PDF embedded image extraction
pip install python-docx         # DOCX / ODT
pip install openpyxl            # XLSX
pip install odfpy               # ODS
pip install python-pptx         # PPTX
pip install pillow pytesseract  # Images OCR (+ Tesseract binary)
pip install trafilatura         # HTML / URL extraction
pip install rarfile             # RAR archives
pip install py7zr               # 7-Zip archives
pip install evtx                # Windows Event Log
pip install defusedxml          # Safe XML parsing (ODT / EVTX)
pip install requests            # URL fetching
pip install duckdb              # DuckDB output format + Explorer
```
""")
    st.divider()
    if st.button("Close", use_container_width=True):
        st.rerun()


# ── Report builder ────────────────────────────────────────────────────────────

def _build_report(results: dict) -> str:
    successes = results["successes"]
    errors    = results["errors"]
    run_ts    = results.get("run_ts", "-")
    total_raw   = sum(r["Tokens Before"] for r in successes) if successes else 0
    total_out   = sum(r["Tokens After"]  for r in successes) if successes else 0
    total_saved = total_raw - total_out
    avg_red     = round((1 - total_out / max(total_raw, 1)) * 100, 1) if total_raw else 0.0

    lines = [
        "# SlimDocs - Extraction Report",
        f"\nGenerated: {run_ts}",
        "\n## Summary",
        f"- Files processed: {len(successes)}",
        f"- Errors: {len(errors)}",
        f"- Tokens before: {total_raw:,}",
        f"- Tokens after:  {total_out:,}",
        f"- Tokens saved:  {total_saved:,}",
        f"- Reduction:     {avg_red}%",
        f"\nOutput: `{results['output_path']}`",
        "\n## Per-file Results",
        "| File | Tokens Before | Tokens After | Saved | Reduction |",
        "|------|--------------|-------------|-------|-----------|",
    ]
    for r in successes:
        lines.append(
            f"| {r['File']} | {r['Tokens Before']:,} | {r['Tokens After']:,} "
            f"| {r['Saved']:,} | {r['Reduction %']}% |"
        )
    if errors:
        lines += ["\n## Errors", "| File | Error |", "|------|-------|"]
        for e in errors:
            name = Path(e["File"]).name if not e["File"].startswith("http") else e["File"][:60]
            lines.append(f"| {name} | {e['Error']} |")
    return "\n".join(lines)


def _on_logs_row_select() -> None:
    """Widget callback for the Logs table: load that run's data for Statistics & Reports.

    ``st.tabs`` has no reliable way to switch the active tab from Python -
    its state-tracking mode (``on_change="rerun"``) reruns the app on every
    tab switch, which both interrupted in-progress extractions and, in
    practice, occasionally reset the active tab on unrelated reruns. So this
    only loads the session's data; the user still clicks the Statistics &
    Reports tab themselves, prompted by the one-shot hint set below.
    """
    sel = st.session_state.get("logs_table")
    if not sel or not sel["selection"]["rows"]:
        return

    logs = st.session_state.logs
    type_filter = st.session_state.get("logs_type_filter", "All")
    search_term = st.session_state.get("logs_search_term", "")
    filtered = logs
    if type_filter != "All":
        filtered = [e for e in filtered if e["Type"] == type_filter]
    if search_term.strip():
        kw = search_term.strip().lower()
        filtered = [
            e for e in filtered
            if kw in e["File"].lower() or kw in e["Message"].lower()
        ]

    idx = sel["selection"]["rows"][0]
    if idx >= len(filtered):
        return
    ts = filtered[idx]["Time"]
    if ts in st.session_state.session_history:
        st.session_state.selected_session_ts = ts
        st.session_state.show_nav_hint = True
        st.toast(f"📊 Session {ts} loaded - open the Statistics & Reports tab to view it")


# ── DuckDB Explorer helpers ───────────────────────────────────────────────────

def _find_snippets(text: str, keyword: str, context: int = 200) -> list:
    """Return one dict per match: surrounding context, line/char position, full record."""
    snippets: list = []
    text_lower = text.lower()
    kw_lower   = keyword.lower()
    kw_len     = max(len(keyword), 1)
    pos, n     = 0, 0
    while True:
        idx = text_lower.find(kw_lower, pos)
        if idx == -1:
            break
        n += 1
        start = max(0, idx - context)
        end   = min(len(text), idx + kw_len + context)

        # Full record: walk back to the non-indented header line owning this match
        line_start = text.rfind("\n", 0, idx) + 1
        cur_start  = line_start
        while cur_start > 0 and text[cur_start: cur_start + 2] == "  ":
            prev_nl   = text.rfind("\n", 0, cur_start - 1)
            cur_start = prev_nl + 1 if prev_nl >= 0 else 0
        line_end = text.find("\n", idx + kw_len)
        if line_end == -1:
            line_end = len(text)
        full_record = text[cur_start:line_end]
        nxt = line_end + 1
        while nxt < len(text):
            nl      = text.find("\n", nxt)
            seg_end = nl if nl != -1 else len(text)
            seg     = text[nxt:seg_end]
            if seg.startswith("  "):
                full_record += "\n" + seg
                nxt = seg_end + 1
            else:
                break

        snippets.append({
            "n":            n,
            "line":         text[:idx].count("\n") + 1,
            "char_pos":     idx,
            "prefix":       "…" if start > 0 else "",
            "text_snippet": text[start:end],
            "suffix":       "…" if end < len(text) else "",
            "full_record":  full_record,
        })
        pos = idx + kw_len
    return snippets


def _render_snippet_html(snippet: dict, keyword: str) -> str:
    """Return highlighted HTML for a single match snippet."""
    import html as _html
    import re as _re

    raw      = snippet["prefix"] + snippet["text_snippet"] + snippet["suffix"]
    escaped  = _html.escape(raw)
    kw_esc   = _re.escape(_html.escape(keyword))
    highlighted = _re.sub(
        kw_esc,
        lambda m: (
            '<mark style="background:#f5c518;color:#000;'
            'padding:0 2px;border-radius:2px;font-weight:bold">'
            + m.group() + "</mark>"
        ),
        escaped,
        flags=_re.IGNORECASE,
    )
    return (
        '<div style="font-family:monospace;font-size:12px;white-space:pre-wrap;'
        "overflow-x:auto;background:#f6f8fa;color:#24292e;"
        "border-left:4px solid #f5c518;padding:10px 14px;"
        'border-radius:4px;line-height:1.6;margin:4px 0">'
        + highlighted + "</div>"
    )


# ── DuckDB Explorer ───────────────────────────────────────────────────────────

def _render_duckdb_explorer(db_path: str):
    """Query Database, SQL editor, and DB visualizer over a slimDocs .duckdb file."""
    try:
        import duckdb as _duckdb
    except ImportError:
        st.error("duckdb not installed - run: pip install duckdb")
        return

    st.caption(f"Database: `{db_path}`")
    tab_search, tab_sql, tab_viz = st.tabs([
        "🔍 Query Database",
        "🗄️ SQL Editor",
        "📊 DB Visualizer",
    ])

    # ═══════════════════════════════════════════════════════════════════════════
    # Query Database
    # ═══════════════════════════════════════════════════════════════════════════
    with tab_search:
        # A plain st.text_input + separate st.button doesn't submit on Enter -
        # Enter only commits the text_input's own value, it doesn't click a
        # nearby button. st.form makes Enter trigger the form's submit button.
        with st.form(key="duckdb_kw_form", border=False):
            kw_col, btn_col = st.columns([5, 1])
            with kw_col:
                keyword = st.text_input(
                    "keyword_input",
                    label_visibility="collapsed",
                    placeholder="Type a keyword to search across all extracted content…",
                )
            with btn_col:
                kw_search = st.form_submit_button("Search", use_container_width=True)

        if kw_search and keyword.strip():
            kw = keyword.strip()
            with st.spinner(f'Searching for "{kw}"…'):
                try:
                    conn = _duckdb.connect(db_path, read_only=True)
                    kw_df = conn.execute(
                        """
                        SELECT
                            doc_id,
                            name,
                            extension,
                            source,
                            size_kb,
                            tokens_before,
                            tokens_after,
                            content,
                            (LENGTH(content)
                             - LENGTH(REPLACE(LOWER(content), LOWER(?), ''))
                            ) / NULLIF(LENGTH(?), 0) AS occurrences
                        FROM documents
                        WHERE content ILIKE '%' || ? || '%'
                        ORDER BY occurrences DESC
                        """,
                        [kw, kw, kw],
                    ).df()
                    conn.close()

                    if kw_df.empty:
                        st.info(f'No documents contain "{kw}".')
                    else:
                        total_occ = int(kw_df["occurrences"].sum())
                        st.success(
                            f'Found **{len(kw_df)}** document(s) with '
                            f'**{total_occ:,}** total occurrence(s) of "{kw}"'
                        )

                        # Summary table (no content column)
                        summary_df = kw_df.drop(columns=["content"])
                        st.dataframe(
                            summary_df,
                            use_container_width=True,
                            hide_index=True,
                            column_config={
                                "doc_id":        st.column_config.NumberColumn("ID",          format="%d"),
                                "name":          st.column_config.TextColumn("File"),
                                "extension":     st.column_config.TextColumn("Ext"),
                                "source":        st.column_config.TextColumn("Source"),
                                "size_kb":       st.column_config.NumberColumn("Size (KB)",   format="%.1f"),
                                "tokens_before": st.column_config.NumberColumn("Tok Before",  format="%d"),
                                "tokens_after":  st.column_config.NumberColumn("Tok After",   format="%d"),
                                "occurrences":   st.column_config.NumberColumn("Occurrences", format="%d"),
                            },
                        )
                        csv = summary_df.to_csv(index=False).encode("utf-8")
                        st.download_button(
                            "📥 Export Search Results",
                            csv,
                            f"search_{kw[:30]}.csv",
                            "text/csv",
                            key="duckdb_export_search",
                        )

                        # Per-file match context panels
                        st.markdown("#### Match Context")
                        st.caption(
                            "Expand a file to browse all highlighted occurrences. "
                            "Click **📋 Full record** on any match to read the complete block."
                        )
                        _MAX_SNIPPETS = 50

                        for _, row in kw_df.iterrows():
                            text     = row["content"] or ""
                            snippets = _find_snippets(text, kw)
                            n_matches = len(snippets)
                            shown     = snippets[:_MAX_SNIPPETS]

                            with st.expander(
                                f"📄 {row['name']}  -  {n_matches:,} match(es)",
                                expanded=False,
                            ):
                                for s in shown:
                                    st.caption(
                                        f"Match {s['n']} of {n_matches} · "
                                        f"Line {s['line']:,} · Char {s['char_pos']:,}"
                                    )
                                    st.markdown(
                                        _render_snippet_html(s, kw),
                                        unsafe_allow_html=True,
                                    )
                                    if s.get("full_record", "").strip():
                                        with st.expander("📋 Full record"):
                                            st.code(s["full_record"], language=None)
                                    st.write("")

                                if n_matches > _MAX_SNIPPETS:
                                    st.info(
                                        f"Showing first {_MAX_SNIPPETS} of {n_matches:,} matches. "
                                        "Use the SQL editor to page through additional results."
                                    )

                except Exception as exc:
                    st.error(f"Search error: {exc}")

    # ═══════════════════════════════════════════════════════════════════════════
    # SQL Editor
    # ═══════════════════════════════════════════════════════════════════════════
    with tab_sql:
        st.subheader("📝 SQL Editor")

        with st.expander("Schema reference & examples", expanded=False):
            st.markdown("""
| Table | Columns |
|-------|---------|
| `documents` | doc_id, name, extension, source, size_kb, tokens_before, tokens_after, tokens_saved, content, extracted_at |

---

**Read file content**
```sql
-- Full extracted text for a specific file
SELECT name, source, content
FROM documents
WHERE name = 'my_file.pdf';

-- Preview first 500 characters of every file
SELECT name, extension,
       LEFT(content, 500) AS preview
FROM documents
WHERE content <> ''
ORDER BY name;
```

**Keyword search across all extracted content**
```sql
-- Find all files whose content contains a keyword
SELECT name, source, extension, tokens_after
FROM documents
WHERE content ILIKE '%error%'
ORDER BY name;

-- Show matching snippet (50 chars before and after keyword)
SELECT name,
       SUBSTRING(content,
           GREATEST(1, STRPOS(LOWER(content), 'error') - 50),
           150
       ) AS snippet
FROM documents
WHERE content ILIKE '%error%';

-- Count occurrences per file
SELECT name,
       (LENGTH(content) - LENGTH(REPLACE(LOWER(content), 'error', '')))
       / LENGTH('error') AS occurrences
FROM documents
WHERE content ILIKE '%error%'
ORDER BY occurrences DESC;
```

**Token savings**
```sql
-- Top 10 files by tokens saved
SELECT name, extension, tokens_before, tokens_after, tokens_saved
FROM documents
ORDER BY tokens_saved DESC LIMIT 10;

-- Summary by file extension
SELECT extension, COUNT(*) AS files,
       SUM(tokens_saved) AS total_saved,
       ROUND(AVG(tokens_saved), 0) AS avg_saved
FROM documents
GROUP BY extension ORDER BY total_saved DESC;
```

**File metadata**
```sql
-- Largest files by size
SELECT name, extension, size_kb
FROM documents
ORDER BY size_kb DESC LIMIT 20;

-- Files extracted most recently
SELECT name, extension, extracted_at
FROM documents
ORDER BY extracted_at DESC LIMIT 20;

-- Files from URLs only
SELECT name, source, tokens_after
FROM documents
WHERE extension = '.url'
ORDER BY name;
```
""")

        query = st.text_area(
            "SQL Query",
            height=130,
            placeholder="SELECT * FROM documents LIMIT 20",
            key="duckdb_sql_input",
        )
        if st.button("▶ Execute", type="primary", key="duckdb_run_sql"):
            if not query.strip():
                st.warning("Enter a SQL query first.")
            else:
                with st.spinner("Running query…"):
                    try:
                        conn = _duckdb.connect(db_path, read_only=True)
                        res_df = conn.execute(query).df()
                        conn.close()
                        st.success(f"{len(res_df)} row(s) returned.")
                        st.dataframe(res_df, use_container_width=True, hide_index=True)
                        csv = res_df.to_csv(index=False).encode("utf-8")
                        st.download_button(
                            "📥 Download Results", csv, "query_results.csv", "text/csv",
                            key="duckdb_export_sql",
                        )
                    except Exception as exc:
                        st.error(f"Query error: {exc}")

    # ═══════════════════════════════════════════════════════════════════════════
    # DB Visualizer
    # ═══════════════════════════════════════════════════════════════════════════
    with tab_viz:
        st.header("🗄️ DB Visualizer")

        _ALL_COLS   = ["doc_id", "name", "extension", "source", "size_kb",
                       "tokens_before", "tokens_after", "tokens_saved",
                       "extracted_at", "content"]
        _LARGE_COLS = {"content"}
        _COL_TYPES  = {
            "doc_id": "INTEGER", "name": "VARCHAR", "extension": "VARCHAR",
            "source": "VARCHAR", "size_kb": "DOUBLE",
            "tokens_before": "INTEGER", "tokens_after": "INTEGER",
            "tokens_saved": "INTEGER", "extracted_at": "TIMESTAMP",
            "content": "TEXT",
        }

        # ── Controls ──────────────────────────────────────────────────────
        ctrl1, ctrl2 = st.columns([4, 1])
        with ctrl1:
            selected_cols = st.multiselect(
                "Columns to display",
                _ALL_COLS,
                default=[c for c in _ALL_COLS if c not in _LARGE_COLS],
                key="dbviz_cols",
            )
        with ctrl2:
            row_limit = st.number_input(
                "Max rows", min_value=10, max_value=50_000,
                value=500, step=100, key="dbviz_limit",
            )

        if not selected_cols:
            st.warning("Select at least one column to display.")
        else:
            # ── Column filters ────────────────────────────────────────────
            with st.expander("🔍 Column Filters", expanded=False):
                st.caption("All filters applied with AND logic (ILIKE).")
                filter_cols = st.multiselect(
                    "Filter columns", selected_cols, key="dbviz_filter_cols",
                )
                filter_values: dict = {}
                if filter_cols:
                    grid = st.columns(min(len(filter_cols), 3))
                    for i, col in enumerate(filter_cols):
                        with grid[i % 3]:
                            val = st.text_input(col, key=f"dbviz_fv_{col}", placeholder="contains…")
                            if val.strip():
                                filter_values[col] = val.strip()

            # ── Sort controls ─────────────────────────────────────────────
            sort_c1, sort_c2 = st.columns([3, 1])
            with sort_c1:
                sort_by = st.selectbox(
                    "Sort by", ["(none)"] + selected_cols, key="dbviz_sort",
                )
            with sort_c2:
                sort_dir = st.radio(
                    "Direction", ["ASC", "DESC"], horizontal=True, key="dbviz_dir",
                )

            # ── Build SQL ─────────────────────────────────────────────────
            col_exprs = [
                f'LEFT(CAST("{c}" AS VARCHAR), 400) AS "{c}"'
                if c in _LARGE_COLS else f'"{c}"'
                for c in selected_cols
            ]
            where_parts, params = [], []
            for col, val in filter_values.items():
                where_parts.append(f'CAST("{col}" AS VARCHAR) ILIKE ?')
                params.append(f"%{val}%")
            where_sql = ("WHERE " + " AND ".join(where_parts)) if where_parts else ""
            order_sql = f'ORDER BY "{sort_by}" {sort_dir}' if sort_by != "(none)" else ""
            sql = (
                f'SELECT {", ".join(col_exprs)} FROM documents '
                f'{where_sql} {order_sql} LIMIT {int(row_limit)}'
            )

            # ── Execute ───────────────────────────────────────────────────
            try:
                conn = _duckdb.connect(db_path, read_only=True)
                df   = conn.execute(sql, params).df()
                total_count = conn.execute(
                    f"SELECT COUNT(*) FROM documents {where_sql}", params
                ).fetchone()[0]
                conn.close()
            except Exception as exc:
                st.error(f"Query error: {exc}")
                st.code(sql, language="sql")
            else:
                status = (
                    f"Showing **{len(df):,}** of **{total_count:,}** rows"
                    + (" (filtered)" if filter_values else "")
                    + (" - `content` truncated to 400 chars"
                       if "content" in selected_cols else "")
                )
                st.caption(status)

                col_cfg: dict = {}
                for c in selected_cols:
                    t = _COL_TYPES.get(c, "")
                    if "TIMESTAMP" in t:
                        col_cfg[c] = st.column_config.DatetimeColumn(c)
                    elif t in ("INTEGER", "DOUBLE"):
                        col_cfg[c] = st.column_config.NumberColumn(c, format="%g")

                st.dataframe(
                    df, use_container_width=True, hide_index=True,
                    column_config=col_cfg or None,
                )

                exp_c1, exp_c2 = st.columns(2)
                with exp_c1:
                    csv = df.to_csv(index=False).encode("utf-8")
                    st.download_button(
                        "📥 Export current view as CSV", csv,
                        "documents_export.csv", "text/csv",
                        use_container_width=True, key="duckdb_export_viz",
                    )
                with exp_c2:
                    st.code(sql, language="sql")


# ── Main app ──────────────────────────────────────────────────────────────────

def main():
    _init_session()
    _inject_keyboard_fix()

    title_col, help_col = st.columns([9, 1])
    with title_col:
        st.title("📄 SlimDocs")
        st.caption(
            "Extract text from any file or URL and reformat it as Markdown, "
            "plain text, chunked Claude-ready JSON, or a queryable DuckDB database - "
            "with token reduction stats."
        )
    with help_col:
        st.write("")
        if st.button("ℹ️ Help", use_container_width=True):
            _show_help()

    # ── Sidebar ───────────────────────────────────────────────────────────────
    with st.sidebar:
        st.header("⚙️ Settings")

        input_mode = st.radio(
            "Input mode",
            ["Single File", "Folder", "URLs"],
            index=["Single File", "Folder", "URLs"].index(st.session_state.input_mode),
            horizontal=True,
        )
        st.session_state.input_mode = input_mode

        # ── Input controls (vary by mode) ─────────────────────────────────
        if input_mode == "URLs":
            urls_text_val = st.text_area(
                "URLs (one per line or comma-separated)",
                value=st.session_state.urls_text,
                height=130,
                placeholder="https://example.com\nhttps://docs.uipath.com/...",
            )
            if urls_text_val != st.session_state.urls_text:
                st.session_state.urls_text = urls_text_val
            recursive = False
        else:
            col_in, col_browse_in = st.columns([5, 1])
            with col_in:
                input_path_val = st.text_input(
                    "Input path",
                    value=st.session_state.input_path,
                    label_visibility="collapsed",
                    placeholder="Paste path or click 📂 / 📄 to browse…",
                )
            with col_browse_in:
                browse_icon = "📂" if input_mode == "Folder" else "📄"
                if st.button(browse_icon, help="Browse"):
                    selected = _pick_folder() if input_mode == "Folder" else _pick_file()
                    if selected:
                        st.session_state.input_path = selected
                        st.rerun()
            if input_path_val != st.session_state.input_path:
                st.session_state.input_path = input_path_val

            recursive = True
            if input_mode == "Folder":
                recursive = st.checkbox("Include subfolders", value=True)

        st.divider()

        # ── Output format ─────────────────────────────────────────────────
        st.subheader("Output Format")
        format_key = st.selectbox(
            "format_select",
            list(_FORMAT_OPTIONS.keys()),
            index=0,
            label_visibility="collapsed",
        )
        fmt = _FORMAT_OPTIONS[format_key]

        st.divider()

        # ── Output destination ────────────────────────────────────────────
        st.subheader("Output Destination")
        col_out, col_browse_out = st.columns([5, 1])
        with col_out:
            output_dir_val = st.text_input(
                "Output directory",
                value=st.session_state.output_dir,
                label_visibility="collapsed",
                placeholder="Output folder path…",
            )
        with col_browse_out:
            if st.button("📂", key="browse_out", help="Browse output folder"):
                selected = _pick_folder()
                if selected:
                    st.session_state.output_dir = selected
                    st.rerun()
        if output_dir_val != st.session_state.output_dir:
            st.session_state.output_dir = output_dir_val

        if fmt == "duckdb":
            zip_output = False
            st.caption("Saves a `.duckdb` file (e.g. `slimDocs_YYYYMMDD_HHMMSS.duckdb`) in the selected folder.")
        else:
            zip_output = st.toggle("Zip output files", value=True)
            st.caption(
                "Saves a single `.zip` archive." if zip_output
                else "Saves files in a timestamped folder."
            )

        st.divider()

        extract_clicked = st.button(
            "🚀 Extract Files",
            use_container_width=True,
            type="primary",
            disabled=st.session_state.processing,
        )

        # ── Extension reference ───────────────────────────────────────────
        with st.expander("Supported file types"):
            for cat, exts in _EXT_CATEGORIES.items():
                st.markdown(f"**{cat}**: {' '.join(sorted(exts))}")

        # ── Progress area (filled during extraction, visible from any tab) ─
        st.divider()
        _sb_progress_slot = st.empty()

    # ── Main tabs ─────────────────────────────────────────────────────────────
    # Deliberately plain (no key/on_change tracking): st.tabs' state-tracking
    # mode reruns the app on every tab switch, which both interrupted
    # in-progress extractions and, in practice, occasionally reset the active
    # tab back to the first one on unrelated reruns (e.g. using the DuckDB
    # Explorer search box). Plain tabs are a pure client-side toggle - no
    # rerun, no interruption, no reset - at the cost of not being able to
    # jump to a tab from Python (see _on_logs_row_select for the fallback).
    if fmt == "duckdb":
        tab1, tab2, tab3, tab4 = st.tabs([
            "📁 File Processing",
            "📊 Statistics & Reports",
            "📋 Logs",
            "🦆 DuckDB Explorer",
        ])
    else:
        tab1, tab2, tab3 = st.tabs([
            "📁 File Processing",
            "📊 Statistics & Reports",
            "📋 Logs",
        ])
        tab4 = None

    # ═══════════════════════════════════════════════════════════════════════════
    # Tab 1 - File Processing
    # ═══════════════════════════════════════════════════════════════════════════
    with tab1:
        if extract_clicked:
            st.session_state.results = None
            out = st.session_state.output_dir.strip()

            if not out:
                st.error("No output directory specified.")
            else:
                output_p = Path(out)
                output_p.mkdir(parents=True, exist_ok=True)

                # ── Resolve inputs ─────────────────────────────────────────
                files:         list = []
                raw_urls:      list = []
                skipped_files: list = []
                scan_root: Path | None = None
                valid = True

                if input_mode == "URLs":
                    raw_urls = [
                        u.strip()
                        for u in re.split(r"[,\n\r]+", st.session_state.urls_text)
                        if u.strip()
                    ]
                    if not raw_urls:
                        st.error("No URLs specified.")
                        valid = False

                elif input_mode == "Single File":
                    inp = st.session_state.input_path.strip()
                    if not inp:
                        st.error("No input path specified.")
                        valid = False
                    elif not Path(inp).is_file():
                        st.error(f"File not found: {inp}")
                        _record_error_only_run(
                            [{"File": inp, "Error": "File not found"}],
                            str(output_p), zip_output,
                        )
                        valid = False
                    else:
                        file_ext = Path(inp).suffix.lower()
                        if file_ext not in _SUPPORTED_EXTENSIONS:
                            st.error(
                                f"Unsupported file type: `{file_ext}`\n\n"
                                f"Supported: {', '.join(sorted(_SUPPORTED_EXTENSIONS))}"
                            )
                            _record_error_only_run(
                                [{"File": inp, "Error": f"Unsupported file type: '{file_ext}'"}],
                                str(output_p), zip_output,
                            )
                            valid = False
                        else:
                            files = [Path(inp)]

                else:  # Folder
                    inp = st.session_state.input_path.strip()
                    if not inp:
                        st.error("No folder path specified.")
                        valid = False
                    elif not Path(inp).is_dir():
                        st.error(f"Folder not found: {inp}")
                        valid = False
                    else:
                        with st.spinner("Scanning folder…"):
                            files = discover_files(Path(inp), recursive=recursive)
                            skipped_files = _discover_skipped_files(Path(inp), recursive=recursive)
                        scan_root = Path(inp)
                        if not files and not skipped_files:
                            st.warning("No files found in the selected folder.")
                            valid = False
                        elif skipped_files:
                            st.warning(
                                f"Skipping {len(skipped_files)} unsupported file(s): "
                                + ", ".join(p.name for p in skipped_files[:10])
                                + ("…" if len(skipped_files) > 10 else "")
                            )

                # ── Run extraction ─────────────────────────────────────────
                if valid and (files or raw_urls):
                    n = len(files) or len(raw_urls)

                    # Progress lives in the sidebar so it stays visible on every tab
                    with _sb_progress_slot.container():
                        st.markdown(f"**⏳ Extracting {n} item(s)…**")
                        progress_bar = st.progress(0)
                        status_text  = st.empty()

                    st.session_state.processing  = True
                    st.session_state.content_map = {}
                    try:
                        if fmt == "duckdb":
                            if raw_urls:
                                result = process_urls_duckdb(
                                    raw_urls, output_p,
                                    progress_bar, status_text,
                                    st.session_state.content_map,
                                )
                            else:
                                result = process_files_duckdb(
                                    files, output_p, scan_root,
                                    progress_bar, status_text,
                                    st.session_state.content_map,
                                )
                        elif raw_urls:
                            result = process_urls(
                                raw_urls, fmt, output_p, zip_output,
                                progress_bar, status_text,
                                st.session_state.content_map,
                            )
                        else:
                            result = process_files(
                                files, fmt, output_p, zip_output,
                                scan_root, progress_bar, status_text,
                                st.session_state.content_map,
                            )
                        if input_mode == "Folder" and skipped_files:
                            _skip_errors = [
                                {"File": str(p), "Error": f"Unsupported file type: '{p.suffix}'"}
                                for p in skipped_files
                            ]
                            result["errors"] = _skip_errors + result["errors"]
                            result["logs"] = (
                                _build_logs([], _skip_errors, result["run_ts"]) + result["logs"]
                            )
                        st.session_state.results = result
                        st.session_state.logs    = result["logs"] + st.session_state.logs
                        st.session_state.session_history[result["run_ts"]] = result
                        st.session_state.selected_session_ts = None
                        # Accumulate session-level totals
                        _suc = result["successes"]
                        _st  = st.session_state.session_totals
                        _st["runs"]     += 1
                        _st["files_ok"] += len(_suc)
                        _st["errors"]   += len(result["errors"])
                        _tb = sum(r["Tokens Before"] for r in _suc)
                        _ta = sum(r["Tokens After"]  for r in _suc)
                        _st["tokens_before"] += _tb
                        _st["tokens_after"]  += _ta
                        _st["tokens_saved"]  += max(0, _tb - _ta)
                    except Exception as e:
                        st.error(f"Unexpected error: {e}")
                    finally:
                        st.session_state.processing = False

                    # Replace progress with compact summary in sidebar
                    _res = st.session_state.results or {}
                    _n_ok  = len(_res.get("successes", []))
                    _n_err = len(_res.get("errors",   []))
                    with _sb_progress_slot.container():
                        if _n_err:
                            st.warning(f"Done - {_n_ok} OK · {_n_err} error(s)")
                        else:
                            st.success(f"✅ Done - {_n_ok} file(s) extracted")

                elif valid and input_mode == "Folder" and skipped_files:
                    # Every file in the folder was unsupported - still record a run so
                    # the skipped files are visible in Logs/Statistics instead of vanishing.
                    _skip_errors = [
                        {"File": str(p), "Error": f"Unsupported file type: '{p.suffix}'"}
                        for p in skipped_files
                    ]
                    _record_error_only_run(_skip_errors, str(output_p), zip_output)
                    with _sb_progress_slot.container():
                        st.warning(f"Done - 0 OK · {len(_skip_errors)} error(s)")

        # ── Results display ────────────────────────────────────────────────
        results = st.session_state.results
        if results:
            successes   = results["successes"]
            errors      = results["errors"]
            total_raw   = sum(r["Tokens Before"] for r in successes) if successes else 0
            total_out   = sum(r["Tokens After"]  for r in successes) if successes else 0
            total_saved = total_raw - total_out
            avg_red     = round((1 - total_out / max(total_raw, 1)) * 100, 1) if total_raw else 0.0

            c1, c2, c3, c4, c5 = st.columns(5)
            c1.metric("Files OK",      len(successes))
            c2.metric("Errors",        len(errors))
            c3.metric("Tokens Before", f"{total_raw:,}")
            c4.metric("Tokens After",  f"{total_out:,}",
                      delta=f"-{total_saved:,} saved", delta_color="inverse")
            c5.metric("Reduction",     f"{avg_red}%")

            if results.get("is_duckdb"):
                output_type = "DuckDB database"
            elif results["zip_output"]:
                output_type = "ZIP archive"
            else:
                output_type = "Folder"
            st.success(f"✅ {output_type} saved to:")
            st.code(results["output_path"])
            st.divider()

            inner_ok, inner_err = st.tabs([
                f"✅ Succeeded ({len(successes)})",
                f"❌ Errors ({len(errors)})",
            ])

            _DISPLAY_COLS = ["File", "Output", "Size (KB)",
                             "Tokens Before", "Tokens After", "Saved", "Reduction %"]
            # URL mode shows an extra URL column
            if successes and "URL" in successes[0]:
                _DISPLAY_COLS = ["File", "URL", "Output", "Size (KB)",
                                 "Tokens Before", "Tokens After", "Saved", "Reduction %"]

            with inner_ok:
                if successes:
                    display_df = pd.DataFrame(
                        [{k: r[k] for k in _DISPLAY_COLS if k in r} for r in successes]
                    )
                    st.caption("Click a row to preview its extracted text and formatted output.")
                    event = st.dataframe(
                        display_df,
                        use_container_width=True,
                        hide_index=True,
                        on_select="rerun",
                        selection_mode="single-row",
                        column_config={
                            "File":          st.column_config.TextColumn("File"),
                            "URL":           st.column_config.TextColumn("URL"),
                            "Output":        st.column_config.TextColumn("Output"),
                            "Size (KB)":     st.column_config.NumberColumn("Size (KB)",     format="%.1f"),
                            "Tokens Before": st.column_config.NumberColumn("Tokens Before", format="%d"),
                            "Tokens After":  st.column_config.NumberColumn("Tokens After",  format="%d"),
                            "Saved":         st.column_config.NumberColumn("Saved",          format="%d"),
                            "Reduction %":   st.column_config.NumberColumn("Reduction %",   format="%.1f%%"),
                        },
                    )
                    if event.selection.rows:
                        st.divider()
                        _show_preview(successes[event.selection.rows[0]])

                    st.write("")
                    csv = display_df.to_csv(index=False).encode("utf-8")
                    st.download_button(
                        "📥 Export results as CSV", csv,
                        "extraction_results.csv", "text/csv",
                    )
                else:
                    st.info("No files were successfully processed.")

            with inner_err:
                if errors:
                    st.dataframe(pd.DataFrame(errors), use_container_width=True, hide_index=True)
                else:
                    st.success("No errors!")

        elif not extract_clicked:
            st.info("👈 Configure the settings in the sidebar, then click **🚀 Extract Files** to begin.")
            with st.expander("ℹ️ Quick guide", expanded=True):
                st.markdown("""
| Setting | Description |
|---------|-------------|
| **Single File** | Pick one file to extract |
| **Folder** | Discover and extract all supported files in a directory tree |
| **URLs** | Paste URLs (comma or newline-separated) - content fetched via trafilatura |
| **md** | Clean Markdown - best for reading (default) |
| **txt** | Plain text - maximum compatibility |
| **json** | Chunked Claude-ready JSON - paste directly into API calls |
| **duckdb** | DuckDB database file - enables full-text search, SQL editor, and table browser |
| **Zip output** | ON = single `.zip`; OFF = timestamped folder *(hidden for DuckDB)* |
""")

    # ═══════════════════════════════════════════════════════════════════════════
    # Tab 2 - Statistics & Reports
    # ═══════════════════════════════════════════════════════════════════════════
    with tab2:
        _sel_ts = st.session_state.selected_session_ts
        if _sel_ts and _sel_ts in st.session_state.session_history:
            results = st.session_state.session_history[_sel_ts]
            _bc1, _bc2 = st.columns([5, 1])
            _bc1.info(f"📌 Viewing session from **{_sel_ts}** (opened from Logs).")
            if _bc2.button("⬅️ Back to latest", use_container_width=True):
                st.session_state.selected_session_ts = None
                st.rerun()
        else:
            results = st.session_state.results
        if not results or (not results["successes"] and not results["errors"]):
            st.info("No extraction results yet - run File Processing first.")
        else:
            try:
                import plotly.express as px
                _has_plotly = True
            except ImportError:
                _has_plotly = False

            # ── Session totals (shown when more than one run has been done) ──
            _stotals = st.session_state.session_totals
            if _stotals["runs"] > 1:
                st.subheader("Session Total")
                _savg = round(
                    (1 - _stotals["tokens_after"] / max(_stotals["tokens_before"], 1)) * 100, 1
                )
                sc1, sc2, sc3, sc4, sc5 = st.columns(5)
                sc1.metric("Runs",          _stotals["runs"])
                sc2.metric("Files OK",      _stotals["files_ok"])
                sc3.metric("Tokens Before", f"{_stotals['tokens_before']:,}")
                sc4.metric("Tokens After",  f"{_stotals['tokens_after']:,}",
                           delta=f"-{_stotals['tokens_saved']:,} saved",
                           delta_color="inverse")
                sc5.metric("Avg Reduction", f"{_savg}%")
                st.divider()
                st.subheader("Current Run" if not _sel_ts else "Selected Run")

            successes   = results["successes"]
            errors      = results["errors"]
            total_raw   = sum(r["Tokens Before"] for r in successes)
            total_out   = sum(r["Tokens After"]  for r in successes)
            total_saved = total_raw - total_out
            avg_red     = round((1 - total_out / max(total_raw, 1)) * 100, 1)

            c1, c2, c3, c4, c5 = st.columns(5)
            c1.metric("Files OK",      len(successes))
            c2.metric("Errors",        len(errors))
            c3.metric("Tokens Before", f"{total_raw:,}")
            c4.metric("Tokens After",  f"{total_out:,}",
                      delta=f"-{total_saved:,} saved", delta_color="inverse")
            c5.metric("Reduction",     f"{avg_red}%")

            st.divider()

            if successes:
                df_s = pd.DataFrame(successes)

                if _has_plotly:
                    col_chart1, col_chart2 = st.columns(2)

                    with col_chart1:
                        st.subheader("Tokens Before vs After")
                        chart_df = df_s[["File", "Tokens Before", "Tokens After"]].copy()
                        chart_df["File"] = chart_df["File"].apply(
                            lambda f: f if len(f) <= 30 else f[:27] + "…"
                        )
                        melted = chart_df.melt(
                            id_vars="File",
                            value_vars=["Tokens Before", "Tokens After"],
                            var_name="Stage", value_name="Tokens",
                        )
                        fig1 = px.bar(
                            melted, x="Tokens", y="File", color="Stage",
                            barmode="group", orientation="h",
                            color_discrete_map={
                                "Tokens Before": "#94a3b8",
                                "Tokens After":  "#22c55e",
                            },
                            height=max(300, len(successes) * 45),
                        )
                        fig1.update_layout(
                            legend=dict(orientation="h", yanchor="bottom", y=1.02),
                            margin=dict(l=0, r=0, t=30, b=0),
                            yaxis=dict(autorange="reversed"),
                        )
                        st.plotly_chart(fig1, use_container_width=True)

                    with col_chart2:
                        st.subheader("Tokens Saved by File Type")
                        df_s["Extension"] = df_s["_ext"]
                        ext_df = df_s.groupby("Extension")["Saved"].sum().reset_index()
                        ext_df = ext_df[ext_df["Saved"] > 0]
                        if not ext_df.empty:
                            fig2 = px.pie(
                                ext_df, values="Saved", names="Extension", hole=0.35, height=350,
                            )
                            fig2.update_traces(textposition="inside", textinfo="percent+label")
                            fig2.update_layout(margin=dict(l=0, r=0, t=30, b=0))
                            st.plotly_chart(fig2, use_container_width=True)
                        else:
                            st.info("No positive token savings to chart.")
                else:
                    st.info("Install plotly (`pip install plotly`) to enable charts.")
                    st.bar_chart(df_s.set_index("File")[["Tokens Before", "Tokens After"]])

                st.divider()
                st.subheader("Full Results Table")
                _DCOLS = ["File", "Output", "Size (KB)", "Tokens Before", "Tokens After", "Saved", "Reduction %"]
                st.dataframe(
                    df_s[[c for c in _DCOLS if c in df_s.columns]],
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "Tokens Before": st.column_config.NumberColumn(format="%d"),
                        "Tokens After":  st.column_config.NumberColumn(format="%d"),
                        "Saved":         st.column_config.NumberColumn(format="%d"),
                        "Reduction %":   st.column_config.NumberColumn(format="%.1f%%"),
                        "Size (KB)":     st.column_config.NumberColumn(format="%.1f"),
                    },
                )
            else:
                st.info("No successfully extracted files in this session.")

            st.divider()
            st.subheader(f"Errors ({len(errors)})")
            if errors:
                df_e = pd.DataFrame(errors)
                st.dataframe(
                    df_e,
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "File":  st.column_config.TextColumn("File",  width="medium"),
                        "Error": st.column_config.TextColumn("Error", width="large"),
                    },
                )
                csv_e = df_e.to_csv(index=False).encode("utf-8")
                st.download_button(
                    "📥 Export errors as CSV", csv_e,
                    "extraction_errors.csv", "text/csv",
                )
            else:
                st.success("No errors in this session!")

            st.divider()
            st.subheader("Download Report")
            report_md = _build_report(results)
            st.markdown(report_md)
            ts_safe = results.get("run_ts", "export").replace(":", "-").replace(" ", "_")
            st.download_button(
                "📥 Download Report (.md)",
                report_md.encode("utf-8"),
                f"token_report_{ts_safe}.md",
                "text/markdown",
            )

    # ═══════════════════════════════════════════════════════════════════════════
    # Tab 3 - Logs
    # ═══════════════════════════════════════════════════════════════════════════
    with tab3:
        logs = st.session_state.logs
        if not logs:
            st.info("No logs yet - run File Processing first.")
        else:
            col_type, col_search = st.columns([2, 4])
            with col_type:
                type_filter = st.selectbox(
                    "Filter by type", ["All", "✅ Success", "❌ Error"], index=0,
                    key="logs_type_filter",
                )
            with col_search:
                search_term = st.text_input(
                    "Search file or message", placeholder="keyword…",
                    key="logs_search_term",
                )

            filtered = logs
            if type_filter != "All":
                filtered = [e for e in filtered if e["Type"] == type_filter]
            if search_term.strip():
                kw = search_term.strip().lower()
                filtered = [
                    e for e in filtered
                    if kw in e["File"].lower() or kw in e["Message"].lower()
                ]

            st.caption(
                f"Showing {len(filtered)} of {len(logs)} entries. "
                "Click a row to load that session's report into Statistics & Reports."
            )
            if st.session_state.show_nav_hint:
                st.session_state.show_nav_hint = False
                st.info("📊 Session loaded - open the **Statistics & Reports** tab above to view it.")
            if filtered:
                log_df = pd.DataFrame(filtered)
                log_event = st.dataframe(
                    log_df,
                    use_container_width=True,
                    hide_index=True,
                    key="logs_table",
                    on_select=_on_logs_row_select,
                    selection_mode="single-row",
                    column_config={
                        "Time":    st.column_config.TextColumn("Time",    width="small"),
                        "Type":    st.column_config.TextColumn("Type",    width="small"),
                        "File":    st.column_config.TextColumn("File",    width="medium"),
                        "Message": st.column_config.TextColumn("Message", width="large"),
                    },
                )
                if (
                    log_event.selection.rows
                    and filtered[log_event.selection.rows[0]]["Time"]
                    not in st.session_state.session_history
                ):
                    st.warning("That session's full data is no longer available.")

                csv = log_df.to_csv(index=False).encode("utf-8")
                st.download_button("📥 Export Logs", csv, "extraction_logs.csv", "text/csv")
            else:
                st.info("No entries match the current filter.")

    # ═══════════════════════════════════════════════════════════════════════════
    # Tab 4 - DuckDB Explorer (only visible when DuckDB format is selected)
    # ═══════════════════════════════════════════════════════════════════════════
    if tab4 is not None:
        with tab4:
            results = st.session_state.results
            if not results or not results.get("is_duckdb"):
                st.info("Run extraction with the **DuckDB** output format to populate the explorer.")
            else:
                _render_duckdb_explorer(results["output_path"])


# ── Entry point ───────────────────────────────────────────────────────────────
main()
